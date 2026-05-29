"""
Prestação de Contas v10 - Associação Maçônica Cidade de Vitória
Exercício 2025 - Dados reais do W011A

Mudanças v10:
- Faixa 121%/17,1% com números BEM maiores
- Patrimônio: removido badge, frase-ponte no subtítulo
- Estrutura de Despesas: paleta verde→azul→cinza (sem vermelho/laranja)
- NOVO slide: Superávit Mensal ao Longo do Ano
- Detalhamentos: template v7 (card esquerdo + tabela direita)
  - Pessoal: 12 lançamentos mês a mês
  - Consumo/Materiais/Serviços/Manutenção/Admin: subcategorias
  - Taxas/Financeiras/Retenções: lançamentos não-zero mês a mês
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ========== PALETA ==========
C_BG_DARK    = RGBColor(0x0A, 0x17, 0x33)
C_NAVY_DEEP  = RGBColor(0x0A, 0x24, 0x63)
C_NAVY       = RGBColor(0x14, 0x3A, 0x87)
C_BLUE       = RGBColor(0x1E, 0x5A, 0xA8)
C_BLUE_MID   = RGBColor(0x2E, 0x7B, 0xC7)
C_BLUE_LIGHT = RGBColor(0x52, 0x99, 0xDC)
C_BLUE_PALE  = RGBColor(0x7F, 0xB5, 0xE3)
C_GRAY_TEXT  = RGBColor(0x3E, 0x56, 0x76)
C_GRAY_MUTED = RGBColor(0x8B, 0x9A, 0xB8)
C_GRAY_BG    = RGBColor(0xF7, 0xF9, 0xFC)
C_GRAY_LINE  = RGBColor(0xE2, 0xE8, 0xF0)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_POSITIVE   = RGBColor(0x22, 0x8B, 0x54)
C_NEGATIVE   = RGBColor(0xC0, 0x3B, 0x3B)

# Destaque
C_AMBER         = RGBColor(0xE8, 0x8B, 0x1A)
C_AMBER_DEEP    = RGBColor(0xC4, 0x6E, 0x0A)
C_AMBER_LIGHT   = RGBColor(0xFD, 0xE9, 0xCC)

# NOVA paleta pra Estrutura de Despesas (verde -> azul -> cinza)
CAT_COLORS = [
    RGBColor(0x1E, 0x73, 0x4A),  # verde escuro (Pessoal - maior)
    RGBColor(0x35, 0x9E, 0x66),  # verde médio (Consumo)
    RGBColor(0x14, 0x3A, 0x87),  # navy (Administrativo)
    RGBColor(0x1E, 0x5A, 0xA8),  # azul (Taxas)
    RGBColor(0x2E, 0x7B, 0xC7),  # azul médio (Materiais)
    RGBColor(0x52, 0x99, 0xDC),  # azul claro (Serviços)
    RGBColor(0x5B, 0x6A, 0x88),  # cinza-azul escuro (Manutenção)
    RGBColor(0x7F, 0x8F, 0xA8),  # cinza-azul (Financeiras)
    RGBColor(0xA5, 0xB0, 0xC2),  # cinza claro (Retenções)
]

FONT = "Calibri"

# ========== DADOS REAIS DO W011A ==========
SALDO_ANTERIOR   = 21040.97
RECEITA_TOTAL    = 116572.72
DESPESA_TOTAL    = 96680.78
SUPERAVIT        = 19891.94
SALDO_FINAL      = 40932.91
MARGEM_STR       = "17,1"
COBERTURA_STR    = "121"
CRESCIMENTO_STR  = "94,5"

MESES = ["Jan","Fev","Mar","Abr","Mai","Jun","Jul","Ago","Set","Out","Nov","Dez"]
MESES_INI = ["J","F","M","A","M","J","J","A","S","O","N","D"]
MESES_EXT = ["Janeiro","Fevereiro","Março","Abril","Maio","Junho",
             "Julho","Agosto","Setembro","Outubro","Novembro","Dezembro"]

RECEITAS_MES = [6099.26, 6423.62, 12336.35, 8612.85, 7965.46, 11864.59,
                10145.25, 9421.49, 8506.36, 10865.71, 10478.22, 13853.56]
DESPESAS_MES = [5725.62, 9746.89, 6883.30, 9285.02, 9830.19, 6822.85,
                6813.12, 6011.24, 10639.85, 9799.34, 9440.33, 5683.03]
SALDO_FIM_MES = [21414.61, 18091.34, 23544.39, 22872.22, 21007.49, 26049.23,
                 29381.36, 32791.61, 30658.12, 31724.49, 32762.38, 40932.91]
SUPERAVIT_MES = [r - d for r, d in zip(RECEITAS_MES, DESPESAS_MES)]
# Meses positivos
MESES_POSITIVOS = sum(1 for v in SUPERAVIT_MES if v > 0)

MEDIA_RECEITA  = RECEITA_TOTAL / 12
MEDIA_DESPESA  = DESPESA_TOTAL / 12
MEDIA_ECONOMIA = MEDIA_RECEITA - MEDIA_DESPESA

RECEITAS_CAT = [
    ("Cotas Ordinárias",      108717.31, 93.3),
    ("Acordos",                 4737.86,  4.1),
    ("Rendimentos",             2235.85,  1.9),
    ("Multas e Juros",           511.94,  0.4),
    ("Outras Receitas",          369.76,  0.3),
]

DESPESAS_CAT = [
    ("Pessoal",               36000.00, 37.2),
    ("Consumo",               19339.17, 20.0),
    ("Administrativo",        11604.23, 12.0),
    ("Taxas e Recolhimentos", 10521.02, 10.9),
    ("Materiais",              6206.33,  6.4),
    ("Serviços",               5162.70,  5.3),
    ("Manutenção",             5109.92,  5.3),
    ("Financeiras",            2164.11,  2.2),
    ("Retenções",               573.30,  0.6),
]

# Mensais por categoria (do W011A)
PESSOAL_TOTAL_MES = [3000.00]*12

CONSUMO_ENERGIA = [397.83, 941.99, 1168.90, 1122.46, 1029.04, 926.82, 784.03, 827.84, 1049.95, 1110.08, 1083.42, 1406.69]
CONSUMO_GAS     = [260.00, 0.00, 650.00, 0.00, 520.00, 0.00, 910.00, 520.00, 260.00, 0.00, 460.00, 460.00]
CONSUMO_INTERNET = [119.90, 119.90, 122.43, 119.90, 119.90, 119.90, 119.90, 119.90, 119.90, 119.90, 119.90, 0.00]
CONSUMO_AGUA    = [119.66, 85.02, 163.82, 163.30, 148.66, 163.82, 223.97, 255.08, 171.88, 462.42, 171.06, 0.00]
CONSUMO_TOTAL_MES = [sum(v) for v in zip(CONSUMO_ENERGIA, CONSUMO_GAS, CONSUMO_INTERNET, CONSUMO_AGUA)]

ADMIN_HONORARIOS = [500.00, 500.00, 500.00, 537.50, 537.50, 537.50, 537.50, 537.50, 537.50, 537.50, 537.50, 549.85]
ADMIN_CARTORIO   = [421.83, 2864.15, 72.61, 1895.79, 0.00, 0.00, 0.00, 0.00, 0.00, 0.00, 0.00, 0.00]
ADMIN_TOTAL_MES = [sum(v) for v in zip(ADMIN_HONORARIOS, ADMIN_CARTORIO)]

TAXAS_IPTU    = [0,0,0,0,0,0,0,0, 3372.38, 3372.35, 3372.35, 0]
TAXAS_ALVARA  = [0,0,0,0, 403.94, 0,0,0,0,0,0,0]
TAXAS_TOTAL_MES = [sum(v) for v in zip(TAXAS_IPTU, TAXAS_ALVARA)]

MAT_LIMPEZA    = [219.35, 390.17, 0, 0, 785.35, 0, 0, 0, 345.11, 0, 0, 0]
MAT_ELETRICO   = [0, 0, 0, 0, 0, 159.60, 0, 48.89, 79.80, 0, 0, 0]
MAT_CONSTRUCAO = [0, 0, 0, 455.00, 0, 0, 0, 0, 0, 0, 0, 0]
MAT_CONSUMO    = [102.87, 695.48, 217.48, 316.20, 75.00, 0, 0, 0, 0, 400.00, 0, 0]
MAT_CONFRAT    = [0, 0, 0, 300.00, 0, 0, 0, 0, 0, 0, 0, 0]
MAT_INCENDIO   = [0, 0, 0, 0, 1083.30, 0, 532.73, 0, 0, 0, 0, 0]
MAT_TOTAL_MES = [sum(v) for v in zip(MAT_LIMPEZA, MAT_ELETRICO, MAT_CONSTRUCAO, MAT_CONSUMO, MAT_CONFRAT, MAT_INCENDIO)]

MANUT_ELEVADOR = [237.00, 237.00, 267.28, 257.00, 257.00, 257.00, 257.00, 257.00, 257.00, 257.00, 257.00, 0.00]
MANUT_SEGURANCA = [208.48, 208.48, 227.84, 208.48, 416.96, 0.00, 208.48, 208.48, 208.48, 208.48, 208.48, 0.00]
MANUT_TOTAL_MES = [sum(v) for v in zip(MANUT_ELEVADOR, MANUT_SEGURANCA)]

SERV_DESINSETIZACAO = [0, 0, 0, 0, 0, 0, 0, 0, 1000.00, 0, 0, 0]
SERV_AR_COND        = [0, 560.00, 0, 0, 0, 1400.00, 0, 0, 0, 0, 0, 0]
SERV_INCENDIO       = [0, 0, 291.06, 724.26, 1187.38, 0, 0, 0, 0, 0, 0, 0]
SERV_TOTAL_MES = [sum(v) for v in zip(SERV_DESINSETIZACAO, SERV_AR_COND, SERV_INCENDIO)]

FIN_TARIFAS = [138.70, 144.70, 201.88, 185.13, 266.16, 174.01, 157.61, 154.65, 155.95, 249.71, 148.72, 184.59]
FIN_IRRF    = [0,0,0,0,0, 2.30, 0,0,0,0,0,0]
FIN_TOTAL_MES = [sum(v) for v in zip(FIN_TARIFAS, FIN_IRRF)]

RET_DAS = [0, 0, 0, 0, 0, 81.90, 81.90, 81.90, 81.90, 81.90, 81.90, 81.90]

# ========== HELPERS ==========
def fmt_brl(v):
    s = f"{v:,.2f}"
    return "R$ " + s.replace(",", "X").replace(".", ",").replace("X", ".")

def fmt_brl_int(v):
    s = f"{v:,.0f}"
    return "R$ " + s.replace(",", ".")

def fmt_pct(v):
    return f"{v:.1f}".replace(".", ",") + "%"

def add_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_rounded_rect(slide, x, y, w, h, fill_color, corner=0.05, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = corner
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color; shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, x, y, w, h, text, size=14, bold=False, color=C_GRAY_TEXT,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def add_multi_run_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    for text, size, bold, color in runs:
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_line(slide, x1, y1, x2, y2, color, weight=1.0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line

def add_footer(slide):
    add_text_box(
        slide, Inches(0.5), Inches(7.1), Inches(12.33), Inches(0.3),
        "ASSOCIAÇÃO MAÇÔNICA CIDADE DE VITÓRIA    •    EXERCÍCIO 2025",
        size=9, color=C_GRAY_MUTED, align=PP_ALIGN.CENTER
    )

def add_slide_header(slide, number, section, title_main, title_accent=None, subtitle=None):
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(10), Inches(0.3),
                 f"{number}    {section}",
                 size=11, bold=True, color=C_BLUE_PALE)
    add_line(slide, Inches(0.5), Inches(0.78), Inches(1.0), Inches(0.78), C_BLUE_MID, weight=1.8)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.33), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = title_main
    r1.font.name = FONT; r1.font.size = Pt(34); r1.font.bold = True
    r1.font.color.rgb = C_NAVY
    if title_accent:
        r2 = p.add_run()
        r2.text = " " + title_accent
        r2.font.name = FONT; r2.font.size = Pt(34); r2.font.bold = True
        r2.font.color.rgb = C_BLUE_MID
    if subtitle:
        add_text_box(slide, Inches(0.5), Inches(1.95), Inches(12.33), Inches(0.4),
                     subtitle, size=14, color=C_GRAY_TEXT)

def style_chart_axes(chart, size=9):
    for ax in [chart.category_axis, chart.value_axis]:
        try:
            ax.tick_labels.font.size = Pt(size)
            ax.tick_labels.font.name = FONT
            ax.tick_labels.font.color.rgb = C_GRAY_TEXT
        except Exception:
            pass


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ============================================================
# SLIDE 1 — CAPA
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = C_BG_DARK
bg.line.fill.background(); bg.shadow.inherit = False

c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-2), Inches(6), Inches(6))
c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x16, 0x2E, 0x5E)
c1.line.fill.background(); c1.shadow.inherit = False

c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(4.5), Inches(4.5), Inches(4.5))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x13, 0x28, 0x54)
c2.line.fill.background(); c2.shadow.inherit = False

# Logo Grupo Service (canto superior direito, acima do círculo decorativo)
import os
# Procura a logo: primeiro na pasta assets/ da skill, depois no caminho de trabalho
_skill_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
_candidates = [
    os.path.join(_skill_dir, "assets", "logo_service_white.png"),
    "/home/claude/logo_service_white.png",
]
logo_path = next((p for p in _candidates if os.path.exists(p)), None)
if logo_path:
    slide.shapes.add_picture(logo_path, Inches(10.4), Inches(0.55),
                             width=Inches(2.4), height=Inches(0.79))

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(0.4),
             "PRESTAÇÃO DE CONTAS", size=14, bold=True, color=C_BLUE_PALE)

tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11), Inches(2.8))
tf = tb.text_frame; tf.word_wrap = True; tf.margin_left = Emu(0)
p1 = tf.paragraphs[0]
r1 = p1.add_run(); r1.text = "Associação Maçônica"
r1.font.name = FONT; r1.font.size = Pt(60); r1.font.bold = True
r1.font.color.rgb = C_WHITE
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "Cidade de Vitória"
r2.font.name = FONT; r2.font.size = Pt(60); r2.font.bold = True
r2.font.color.rgb = C_BLUE_LIGHT

add_line(slide, Inches(0.8), Inches(5.2), Inches(2.0), Inches(5.2), C_AMBER, weight=3)
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(10), Inches(0.4),
             "Exercício 2025", size=20, bold=True, color=C_WHITE)
add_text_box(slide, Inches(0.8), Inches(5.98), Inches(10), Inches(0.4),
             "Janeiro a Dezembro  •  Apresentação em Assembleia",
             size=13, color=C_BLUE_PALE)


# ============================================================
# SLIDE 2 — VISÃO GERAL (faixa 121% e 17,1% BEM destacados)
# ============================================================
slide = prs.slides.add_slide(BLANK)
add_slide_header(slide, "01", "VISÃO GERAL", "Resultado consolidado", "do exercício")

card_h = Inches(1.55)
gap = Inches(0.2)
left_margin = Inches(0.5)
row1_y = Inches(2.2)
row2_y = Inches(3.95)
card_w_3 = Inches((13.333 - 1.0 - 0.4) / 3)

kpis_row1 = [
    ("SALDO INICIAL",  fmt_brl(SALDO_ANTERIOR), "Em 01/01/2025", C_BLUE),
    ("RECEITA TOTAL",  fmt_brl(RECEITA_TOTAL),  "12 meses",      C_NAVY),
    ("DESPESA TOTAL",  fmt_brl(DESPESA_TOTAL),  "12 meses",      C_BLUE_MID),
]
x = left_margin
for label, value, sub, color in kpis_row1:
    add_rounded_rect(slide, x, row1_y, card_w_3, card_h, color, corner=0.08)
    add_text_box(slide, x + Inches(0.25), row1_y + Inches(0.18), card_w_3 - Inches(0.5), Inches(0.3),
                 label, size=10, bold=True, color=C_WHITE)
    add_text_box(slide, x + Inches(0.25), row1_y + Inches(0.5), card_w_3 - Inches(0.5), Inches(0.7),
                 value, size=26, bold=True, color=C_WHITE)
    add_text_box(slide, x + Inches(0.25), row1_y + Inches(1.2), card_w_3 - Inches(0.5), Inches(0.3),
                 sub, size=11, color=C_BLUE_PALE)
    x += card_w_3 + gap

card_w_2 = Inches((13.333 - 1.0 - 0.2) / 2)
add_rounded_rect(slide, left_margin, row2_y, card_w_2, card_h, C_POSITIVE, corner=0.08)
add_text_box(slide, left_margin + Inches(0.25), row2_y + Inches(0.18), card_w_2 - Inches(0.5), Inches(0.3),
             "SUPERÁVIT DO EXERCÍCIO", size=10, bold=True, color=C_WHITE)
add_text_box(slide, left_margin + Inches(0.25), row2_y + Inches(0.5), card_w_2 - Inches(0.5), Inches(0.7),
             fmt_brl(SUPERAVIT), size=30, bold=True, color=C_WHITE)
add_text_box(slide, left_margin + Inches(0.25), row2_y + Inches(1.2), card_w_2 - Inches(0.5), Inches(0.3),
             f"Margem de {MARGEM_STR}% sobre receitas", size=11, color=C_WHITE)

x2 = left_margin + card_w_2 + gap
add_rounded_rect(slide, x2, row2_y, card_w_2, card_h, C_NAVY_DEEP, corner=0.08)
add_text_box(slide, x2 + Inches(0.25), row2_y + Inches(0.18), card_w_2 - Inches(0.5), Inches(0.3),
             "SALDO FINAL", size=10, bold=True, color=C_WHITE)
add_text_box(slide, x2 + Inches(0.25), row2_y + Inches(0.5), card_w_2 - Inches(0.5), Inches(0.7),
             fmt_brl(SALDO_FINAL), size=30, bold=True, color=C_WHITE)
add_text_box(slide, x2 + Inches(0.25), row2_y + Inches(1.2), card_w_2 - Inches(0.5), Inches(0.3),
             "Reserva em caixa em 31/12/2025", size=11, color=C_BLUE_PALE)

# ========== FAIXA DE DESTAQUE GIGANTE ==========
diag_y = Inches(5.65)
add_rounded_rect(slide, left_margin, diag_y, Inches(12.33), Inches(1.35), C_AMBER, corner=0.06)

# Dividir em 3 "pilares" visuais dentro da faixa
pilar_w = Inches(12.33 / 3)

# Pilar 1: 121%
add_text_box(slide, left_margin + Inches(0.2), diag_y + Inches(0.12),
             pilar_w - Inches(0.4), Inches(0.3),
             "RECEITA COBRIU", size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, left_margin + Inches(0.2), diag_y + Inches(0.38),
             pilar_w - Inches(0.4), Inches(0.7),
             COBERTURA_STR + "%", size=50, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, left_margin + Inches(0.2), diag_y + Inches(1.02),
             pilar_w - Inches(0.4), Inches(0.3),
             "da despesa anual", size=10, color=C_WHITE, align=PP_ALIGN.CENTER)

# Divisor 1
add_line(slide,
         left_margin + pilar_w, diag_y + Inches(0.25),
         left_margin + pilar_w, diag_y + Inches(1.1),
         C_WHITE, weight=0.8)

# Pilar 2: 17,1%
add_text_box(slide, left_margin + pilar_w + Inches(0.2), diag_y + Inches(0.12),
             pilar_w - Inches(0.4), Inches(0.3),
             "SUPERÁVIT DE", size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, left_margin + pilar_w + Inches(0.2), diag_y + Inches(0.38),
             pilar_w - Inches(0.4), Inches(0.7),
             MARGEM_STR + "%", size=50, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, left_margin + pilar_w + Inches(0.2), diag_y + Inches(1.02),
             pilar_w - Inches(0.4), Inches(0.3),
             "sobre as receitas", size=10, color=C_WHITE, align=PP_ALIGN.CENTER)

# Divisor 2
add_line(slide,
         left_margin + pilar_w * 2, diag_y + Inches(0.25),
         left_margin + pilar_w * 2, diag_y + Inches(1.1),
         C_WHITE, weight=0.8)

# Pilar 3: saldo dobrou
add_text_box(slide, left_margin + pilar_w * 2 + Inches(0.2), diag_y + Inches(0.12),
             pilar_w - Inches(0.4), Inches(0.3),
             "SALDO EM CAIXA", size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, left_margin + pilar_w * 2 + Inches(0.2), diag_y + Inches(0.38),
             pilar_w - Inches(0.4), Inches(0.7),
             "QUASE DOBROU", size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, left_margin + pilar_w * 2 + Inches(0.2), diag_y + Inches(1.02),
             pilar_w - Inches(0.4), Inches(0.3),
             "de R$ 21 mil para R$ 40 mil", size=10, color=C_WHITE, align=PP_ALIGN.CENTER)

add_footer(slide)


# ============================================================
# SLIDE 3 — EVOLUÇÃO MENSAL (Receita x Despesa x Saldo)
# ============================================================
slide = prs.slides.add_slide(BLANK)
add_slide_header(slide, "02", "EVOLUÇÃO MENSAL", "Receita, despesa", "e saldo acumulado",
                 subtitle="Movimentação mensal ao longo do exercício")

chart_data = CategoryChartData()
chart_data.categories = MESES
chart_data.add_series("Receitas",        RECEITAS_MES)
chart_data.add_series("Despesas",        DESPESAS_MES)
chart_data.add_series("Saldo Acumulado", SALDO_FIM_MES)

chart_shape = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, Inches(0.5), Inches(2.5),
    Inches(12.33), Inches(3.5), chart_data
)
chart = chart_shape.chart
chart.has_title = False
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.TOP
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(11); chart.legend.font.name = FONT

plot = chart.plots[0]
for i, col in enumerate([C_POSITIVE, C_NEGATIVE, C_BLUE_MID]):
    s = plot.series[i]
    s.format.line.color.rgb = col
    s.format.line.width = Pt(2.5)
style_chart_axes(chart)

cards_y = Inches(6.15)
card_w = Inches((13.333 - 1.0 - 0.4) / 3)
card_h = Inches(0.85)
medias = [
    ("MÉDIA MENSAL DE RECEITA",  fmt_brl(MEDIA_RECEITA),  C_POSITIVE),
    ("MÉDIA MENSAL DE DESPESA",  fmt_brl(MEDIA_DESPESA),  C_NEGATIVE),
    ("ECONOMIA MÉDIA MENSAL",    fmt_brl(MEDIA_ECONOMIA), C_BLUE),
]
x = Inches(0.5)
for label, value, color in medias:
    add_rounded_rect(slide, x, cards_y, card_w, card_h, color, corner=0.12)
    add_text_box(slide, x + Inches(0.2), cards_y + Inches(0.08), card_w - Inches(0.4), Inches(0.3),
                 label, size=9, bold=True, color=C_WHITE)
    add_text_box(slide, x + Inches(0.2), cards_y + Inches(0.35), card_w - Inches(0.4), Inches(0.5),
                 value, size=18, bold=True, color=C_WHITE)
    x += card_w + Inches(0.2)
add_footer(slide)


# ============================================================
# SLIDE 4 — PATRIMÔNIO (frase-ponte + card +94,5%)
# ============================================================
slide = prs.slides.add_slide(BLANK)
add_slide_header(slide, "03", "PATRIMÔNIO", "Evolução do saldo", "em caixa",
                 subtitle="O destaque do ano foi para o patrimônio em caixa, que cresceu")

chart_data2 = CategoryChartData()
chart_data2.categories = MESES
chart_data2.add_series("Saldo em Caixa", SALDO_FIM_MES)
chart_shape2 = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, Inches(0.5), Inches(2.6),
    Inches(8.3), Inches(3.9), chart_data2
)
chart2 = chart_shape2.chart
chart2.has_title = False; chart2.has_legend = False
s = chart2.plots[0].series[0]
s.format.line.color.rgb = C_BLUE
s.format.line.width = Pt(3.5)
style_chart_axes(chart2)

# Cards laterais
card_x = Inches(9.1); card_w = Inches(3.7)

# Card Crescimento — GRANDE E EM LARANJA (destaque principal)
add_rounded_rect(slide, card_x, Inches(2.6), card_w, Inches(1.75), C_AMBER, corner=0.08)
add_text_box(slide, card_x + Inches(0.25), Inches(2.8), card_w - Inches(0.5), Inches(0.3),
             "CRESCIMENTO NO ANO", size=11, bold=True, color=C_WHITE)
add_text_box(slide, card_x + Inches(0.25), Inches(3.18), card_w - Inches(0.5), Inches(1.1),
             f"+{CRESCIMENTO_STR}%", size=48, bold=True, color=C_WHITE)

# Saldo Inicial
add_rounded_rect(slide, card_x, Inches(4.5), card_w, Inches(0.65), C_BLUE_LIGHT, corner=0.08)
add_text_box(slide, card_x + Inches(0.25), Inches(4.58), card_w - Inches(0.5), Inches(0.22),
             "SALDO INICIAL (JAN)", size=9, bold=True, color=C_WHITE)
add_text_box(slide, card_x + Inches(0.25), Inches(4.8), card_w - Inches(0.5), Inches(0.35),
             fmt_brl(SALDO_ANTERIOR), size=15, bold=True, color=C_WHITE)

# Saldo Final
add_rounded_rect(slide, card_x, Inches(5.25), card_w, Inches(0.65), C_BLUE_MID, corner=0.08)
add_text_box(slide, card_x + Inches(0.25), Inches(5.33), card_w - Inches(0.5), Inches(0.22),
             "SALDO FINAL (DEZ)", size=9, bold=True, color=C_WHITE)
add_text_box(slide, card_x + Inches(0.25), Inches(5.55), card_w - Inches(0.5), Inches(0.35),
             fmt_brl(SALDO_FINAL), size=15, bold=True, color=C_WHITE)

# Diferença
add_rounded_rect(slide, card_x, Inches(6.0), card_w, Inches(0.65), C_POSITIVE, corner=0.08)
add_text_box(slide, card_x + Inches(0.25), Inches(6.08), card_w - Inches(0.5), Inches(0.22),
             "DIFERENÇA / SUPERÁVIT", size=9, bold=True, color=C_WHITE)
add_text_box(slide, card_x + Inches(0.25), Inches(6.3), card_w - Inches(0.5), Inches(0.35),
             fmt_brl(SUPERAVIT), size=15, bold=True, color=C_WHITE)

add_footer(slide)


# ============================================================
# SLIDE 5 — SUPERÁVIT MENSAL (novo, baseado na v7)
# ============================================================
slide = prs.slides.add_slide(BLANK)
add_slide_header(slide, "04", "RESULTADO LÍQUIDO", "Superávit mensal", "ao longo do ano",
                 subtitle="Resultado de cada mês (receitas menos despesas)")

# Gráfico de barras (coluna)
chart_data_sup = CategoryChartData()
chart_data_sup.categories = MESES
chart_data_sup.add_series("Superávit", SUPERAVIT_MES)

chart_shape_sup = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(2.5),
    Inches(8.3), Inches(4.2), chart_data_sup
)
chart_sup = chart_shape_sup.chart
chart_sup.has_title = False; chart_sup.has_legend = False
plot_sup = chart_sup.plots[0]
s_sup = plot_sup.series[0]
s_sup.format.fill.solid()
s_sup.format.fill.fore_color.rgb = C_BLUE
s_sup.format.line.fill.background()
style_chart_axes(chart_sup)

# Cards laterais
card_x = Inches(9.1); card_w = Inches(3.7)

# Meses positivos
add_rounded_rect(slide, card_x, Inches(2.6), card_w, Inches(1.9), C_POSITIVE, corner=0.08)
add_text_box(slide, card_x + Inches(0.25), Inches(2.8), card_w - Inches(0.5), Inches(0.3),
             "MESES POSITIVOS", size=11, bold=True, color=C_WHITE)
add_multi_run_text(
    slide, card_x + Inches(0.25), Inches(3.2), card_w - Inches(0.5), Inches(1.3),
    [(f"{MESES_POSITIVOS}", 80, True, C_WHITE),
     (" de 12", 22, False, C_WHITE)],
    anchor=MSO_ANCHOR.MIDDLE
)

# Superávit anual
add_rounded_rect(slide, card_x, Inches(4.65), card_w, Inches(2.0), C_NAVY_DEEP, corner=0.08)
add_text_box(slide, card_x + Inches(0.25), Inches(4.85), card_w - Inches(0.5), Inches(0.3),
             "SUPERÁVIT ANUAL", size=11, bold=True, color=C_WHITE)
add_text_box(slide, card_x + Inches(0.25), Inches(5.25), card_w - Inches(0.5), Inches(0.9),
             fmt_brl(SUPERAVIT), size=26, bold=True, color=C_WHITE)
add_text_box(slide, card_x + Inches(0.25), Inches(6.25), card_w - Inches(0.5), Inches(0.3),
             "Resultado positivo consolidado.", size=10, color=C_BLUE_PALE)

add_footer(slide)


# ============================================================
# SLIDE 6 — ORIGEM DA RECEITA
# ============================================================
slide = prs.slides.add_slide(BLANK)
add_slide_header(slide, "05", "ORIGEM DA RECEITA", "De onde veio", "a arrecadação",
                 subtitle=f"Total: {fmt_brl(RECEITA_TOTAL)}    •    Exercício 2025")

add_text_box(slide, Inches(0.5), Inches(2.7), Inches(4.5), Inches(0.5),
             fmt_brl_int(RECEITA_TOTAL), size=40, bold=True, color=C_NAVY)
add_text_box(slide, Inches(0.5), Inches(3.55), Inches(4.5), Inches(0.3),
             "RECEITA TOTAL", size=11, bold=True, color=C_GRAY_MUTED)

list_x = Inches(5.2); list_w = Inches(7.6)
row_h = Inches(0.55); list_y = Inches(2.55)
for i, (fonte, val, pct) in enumerate(RECEITAS_CAT):
    y = list_y + i * (row_h + Inches(0.1))
    color = C_POSITIVE if i == 0 else C_BLUE
    add_rounded_rect(slide, list_x, y, list_w, row_h, color, corner=0.12)
    add_text_box(slide, list_x + Inches(0.25), y + Inches(0.13), Inches(3), Inches(0.28),
                 fonte, size=13, bold=True, color=C_WHITE)
    add_text_box(slide, list_x + Inches(3.3), y + Inches(0.15), Inches(2.5), Inches(0.28),
                 fmt_brl(val), size=11, color=C_WHITE)
    add_text_box(slide, list_x + Inches(5.9), y + Inches(0.11), Inches(1.5), Inches(0.3),
                 fmt_pct(pct), size=15, bold=True, color=C_WHITE, align=PP_ALIGN.RIGHT)

diag_y = Inches(6.3)
add_rounded_rect(slide, Inches(0.5), diag_y, Inches(12.33), Inches(0.6), C_AMBER_LIGHT, corner=0.15)
add_multi_run_text(
    slide, Inches(0.8), diag_y + Inches(0.15), Inches(12), Inches(0.3),
    [
        ("Cotas ordinárias respondem por ", 13, False, C_AMBER_DEEP),
        ("93,3%", 14, True, C_AMBER_DEEP),
        (" da arrecadação, base de receita estável e previsível.", 13, False, C_AMBER_DEEP),
    ]
)
add_footer(slide)


# ============================================================
# SLIDE 7 — ESTRUTURA DE DESPESAS (nova paleta verde→azul→cinza)
# ============================================================
slide = prs.slides.add_slide(BLANK)
add_slide_header(slide, "06", "ESTRUTURA DE DESPESAS", "Como a despesa", "foi distribuída",
                 subtitle=f"Total: {fmt_brl(DESPESA_TOTAL)}    •    9 categorias principais")

add_text_box(slide, Inches(0.5), Inches(2.7), Inches(4.5), Inches(0.5),
             fmt_brl_int(DESPESA_TOTAL), size=40, bold=True, color=C_NAVY)
add_text_box(slide, Inches(0.5), Inches(3.55), Inches(4.5), Inches(0.3),
             "DESPESA TOTAL", size=11, bold=True, color=C_GRAY_MUTED)
# Legenda nova
add_text_box(slide, Inches(0.5), Inches(4.0), Inches(4.5), Inches(0.3),
             "Tons verdes: maior peso", size=10, color=C_POSITIVE)
add_text_box(slide, Inches(0.5), Inches(4.25), Inches(4.5), Inches(0.3),
             "Tons azuis e cinzas: peso intermediário e menor", size=10, color=C_GRAY_TEXT)

list_x = Inches(5.2); list_w = Inches(7.6)
row_h = Inches(0.38); list_y = Inches(2.5)
for i, (cat, val, pct) in enumerate(DESPESAS_CAT):
    y = list_y + i * (row_h + Inches(0.06))
    color = CAT_COLORS[i]
    add_rounded_rect(slide, list_x, y, list_w, row_h, color, corner=0.15)
    add_text_box(slide, list_x + Inches(0.25), y + Inches(0.07), Inches(3), Inches(0.25),
                 cat, size=11, bold=True, color=C_WHITE)
    add_text_box(slide, list_x + Inches(3.3), y + Inches(0.08), Inches(2.5), Inches(0.25),
                 fmt_brl(val), size=10, color=C_WHITE)
    add_text_box(slide, list_x + Inches(5.9), y + Inches(0.06), Inches(1.5), Inches(0.28),
                 fmt_pct(pct), size=13, bold=True, color=C_WHITE, align=PP_ALIGN.RIGHT)
add_footer(slide)


# ============================================================
# TEMPLATE DE DETALHAMENTO (estilo v7 — card esquerdo + tabela direita)
# ============================================================
def build_detail_v7(numero, categoria_titulo_1, categoria_titulo_2,
                    total, pct_total, descricao, serie_mensal, lancamentos,
                    icon_shape=MSO_SHAPE.RECTANGLE):
    """
    categoria_titulo_1/2: "Despesas com" / "pessoal" (o segundo fica em cor destaque)
    descricao: texto de 1-2 linhas
    serie_mensal: lista de 12 valores (mostrado como barras J F M A M J J A S O N D)
    lancamentos: lista de tuplas (descrição, valor) — pode ter qualquer número de itens
    """
    slide = prs.slides.add_slide(BLANK)
    num_txt = numero.zfill(2)

    # Header institucional no topo
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(10), Inches(0.3),
                 f"{num_txt}    DETALHAMENTO DE DESPESAS", size=11, bold=True, color=C_BLUE_PALE)
    add_line(slide, Inches(0.5), Inches(0.78), Inches(1.0), Inches(0.78), C_BLUE_MID, weight=1.8)

    # Título
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.33), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = categoria_titulo_1
    r1.font.name = FONT; r1.font.size = Pt(34); r1.font.bold = True
    r1.font.color.rgb = C_NAVY
    if categoria_titulo_2:
        r2 = p.add_run()
        r2.text = " " + categoria_titulo_2
        r2.font.name = FONT; r2.font.size = Pt(34); r2.font.bold = True
        r2.font.color.rgb = C_BLUE_MID

    # ========== CARD ESQUERDO (AZUL ESCURO) ==========
    card_x = Inches(0.5); card_y = Inches(2.15)
    card_w = Inches(5.2); card_h = Inches(4.75)
    add_rounded_rect(slide, card_x, card_y, card_w, card_h, C_NAVY, corner=0.04)

    # Ícone (círculo decorativo)
    icon_x = card_x + Inches(0.35); icon_y = card_y + Inches(0.35)
    icon_s = Inches(0.5)
    icon = slide.shapes.add_shape(icon_shape, icon_x, icon_y, icon_s, icon_s)
    icon.fill.solid(); icon.fill.fore_color.rgb = C_BLUE_MID
    icon.line.fill.background(); icon.shadow.inherit = False

    # "TOTAL DO GRUPO"
    add_text_box(slide, card_x + Inches(0.35), card_y + Inches(1.1), card_w - Inches(0.7), Inches(0.3),
                 "TOTAL DO GRUPO", size=11, bold=True, color=C_BLUE_PALE)

    # Valor grande
    add_text_box(slide, card_x + Inches(0.35), card_y + Inches(1.45), card_w - Inches(0.7), Inches(0.8),
                 fmt_brl(total), size=34, bold=True, color=C_WHITE)

    # Percentual
    add_text_box(slide, card_x + Inches(0.35), card_y + Inches(2.28), card_w - Inches(0.7), Inches(0.3),
                 f"{fmt_pct(pct_total)} DA DESPESA ANUAL", size=10, bold=True, color=C_BLUE_PALE)

    # Linha decorativa
    add_line(slide, card_x + Inches(0.35), card_y + Inches(2.65),
             card_x + Inches(1.35), card_y + Inches(2.65), C_BLUE_LIGHT, weight=2)

    # Descrição
    add_text_box(slide, card_x + Inches(0.35), card_y + Inches(2.85), card_w - Inches(0.7), Inches(0.8),
                 descricao, size=11, color=C_WHITE)

    # "DISTRIBUIÇÃO MENSAL"
    add_text_box(slide, card_x + Inches(0.35), card_y + Inches(3.75), card_w - Inches(0.7), Inches(0.3),
                 "DISTRIBUIÇÃO MENSAL", size=10, bold=True, color=C_BLUE_PALE)

    # Gráfico de barras embutido no card
    chart_w = card_w - Inches(0.7)
    chart_h = Inches(0.85)
    chart_data_det = CategoryChartData()
    chart_data_det.categories = MESES_INI
    chart_data_det.add_series("Mensal", serie_mensal)

    chart_shape_det = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        card_x + Inches(0.35), card_y + Inches(4.05),
        chart_w, chart_h, chart_data_det
    )
    chart_det = chart_shape_det.chart
    chart_det.has_title = False
    chart_det.has_legend = False
    plot_det = chart_det.plots[0]
    s_det = plot_det.series[0]
    s_det.format.fill.solid()
    s_det.format.fill.fore_color.rgb = C_WHITE
    s_det.format.line.fill.background()

    # Esconder eixos no mini gráfico
    try:
        chart_det.value_axis.visible = False
    except Exception:
        pass
    try:
        chart_det.category_axis.tick_labels.font.size = Pt(7)
        chart_det.category_axis.tick_labels.font.color.rgb = C_BLUE_PALE
    except Exception:
        pass

    # ========== TABELA DIREITA (LANÇAMENTOS) ==========
    tab_x = Inches(6.0); tab_y = Inches(2.15)
    tab_w = Inches(6.83)

    # Cabeçalho
    header_h = Inches(0.45)
    add_rect(slide, tab_x, tab_y, tab_w, header_h, C_NAVY_DEEP)
    add_text_box(slide, tab_x + Inches(0.3), tab_y + Inches(0.1), Inches(4), Inches(0.28),
                 "LANÇAMENTOS DO EXERCÍCIO", size=10, bold=True, color=C_WHITE)
    add_text_box(slide, tab_x + tab_w - Inches(1.5), tab_y + Inches(0.1), Inches(1.2), Inches(0.28),
                 "VALOR", size=10, bold=True, color=C_WHITE, align=PP_ALIGN.RIGHT)

    # Linhas: altura adaptativa
    available_h_emu = Inches(4.30).emu  # espaço até o total
    n = len(lancamentos)
    total_reserve_emu = Inches(0.55).emu
    line_h_emu = max(int((available_h_emu - total_reserve_emu) / n), Inches(0.26).emu)
    line_h = Emu(line_h_emu)

    current_y = tab_y + header_h
    font_desc = 10 if n <= 12 else (9 if n <= 16 else 8)
    font_val  = 11 if n <= 12 else (10 if n <= 16 else 9)

    for i, (desc, val) in enumerate(lancamentos):
        # fundo alternado branco/cinza muito claro
        if i % 2 == 1:
            add_rect(slide, tab_x, current_y, tab_w, line_h, C_GRAY_BG)
        # Descrição
        add_text_box(slide,
                     tab_x + Inches(0.3),
                     current_y + Emu(int(line_h_emu/2) - int(Inches(0.13).emu)),
                     Inches(4.5), Inches(0.26),
                     desc, size=font_desc, color=C_GRAY_TEXT)
        # Valor
        add_text_box(slide,
                     tab_x + tab_w - Inches(1.8),
                     current_y + Emu(int(line_h_emu/2) - int(Inches(0.13).emu)),
                     Inches(1.5), Inches(0.26),
                     fmt_brl(val), size=font_val, bold=True, color=C_NAVY, align=PP_ALIGN.RIGHT)
        current_y = current_y + line_h

    # Faixa total (azul escuro)
    total_h = Inches(0.45)
    add_rect(slide, tab_x, current_y + Inches(0.05), tab_w, total_h, C_NAVY)
    add_text_box(slide, tab_x + Inches(0.3), current_y + Inches(0.15), Inches(4), Inches(0.28),
                 f"TOTAL {categoria_titulo_1.upper()} {categoria_titulo_2.upper() if categoria_titulo_2 else ''}".strip(),
                 size=11, bold=True, color=C_WHITE)
    add_text_box(slide, tab_x + tab_w - Inches(2.0), current_y + Inches(0.13), Inches(1.7), Inches(0.3),
                 fmt_brl(total), size=13, bold=True, color=C_WHITE, align=PP_ALIGN.RIGHT)

    add_footer(slide)
    return slide


# ========== SLIDE DETALHAMENTOS ==========

# 8 — Pessoal: 12 lançamentos mês a mês
pessoal_lanc = [(f"Salários Contratados — {MESES_EXT[i]}/2025", 3000.00) for i in range(12)]
build_detail_v7(
    "07", "Despesas com", "pessoal",
    36000.00, 37.2,
    "Folha mensal dos colaboradores contratados pela Associação. Pagamento fixo de R$ 3.000,00 mensais, sem variações ao longo de todo o exercício.",
    PESSOAL_TOTAL_MES,
    pessoal_lanc
)

# 9 — Consumo: subcategorias agrupadas
consumo_lanc = [
    ("Energia Elétrica (fornecimento mensal)", 11849.05),
    ("Gás (abastecimento)",                     4040.00),
    ("Água e Esgoto (fornecimento mensal)",     2128.69),
    ("Internet (serviço fixo mensal)",          1321.43),
]
build_detail_v7(
    "08", "Despesas com", "consumo",
    19339.17, 20.0,
    "Serviços essenciais de operação da sede. Energia elétrica é o item de maior peso (61% da categoria), seguida por gás, água e internet.",
    CONSUMO_TOTAL_MES,
    consumo_lanc
)

# 10 — Administrativo
admin_lanc = [
    ("Honorários Administrativos (administração mensal)", 6349.85),
    ("Despesas com Cartório (registros pontuais)",        5254.38),
]
build_detail_v7(
    "09", "Despesas", "administrativas",
    11604.23, 12.0,
    "Serviços de gestão e registros cartoriais. Honorários são pagos mensalmente. Cartório concentrado em Fev e Abr.",
    ADMIN_TOTAL_MES,
    admin_lanc
)

# 11 — Taxas: mês a mês completo (poucos lançamentos reais)
taxas_lanc = [
    ("Alvará — Maio/2025",           403.94),
    ("IPTU — parcela Setembro/2025", 3372.38),
    ("IPTU — parcela Outubro/2025",  3372.35),
    ("IPTU — parcela Novembro/2025", 3372.35),
]
build_detail_v7(
    "10", "Taxas", "e recolhimentos",
    10521.02, 10.9,
    "IPTU pago em 3 parcelas iguais entre Set e Nov. Alvará em parcela única em Mai/2025.",
    TAXAS_TOTAL_MES,
    taxas_lanc
)

# 12 — Materiais: subcategorias
materiais_lanc = [
    ("Material de Uso e Consumo (itens operacionais)",      1807.03),
    ("Material de Limpeza (produtos, descartáveis, panos)", 1739.98),
    ("Material Sistema de Incêndio (recarga, mangueiras)",  1616.03),
    ("Material de Construção e Ferramentas",                 455.00),
    ("Material de Confraternização (evento institucional)",  300.00),
    ("Material Elétrico (lâmpadas, fusíveis, reposições)",   288.29),
]
build_detail_v7(
    "11", "Aquisição de", "materiais",
    6206.33, 6.4,
    "Insumos para manutenção, limpeza, operação da sede, evento de confraternização e reposição do sistema de segurança.",
    MAT_TOTAL_MES,
    materiais_lanc
)

# 13 — Serviços: subcategorias
servicos_lanc = [
    ("Sistema de Incêndio (manutenção preventiva)",  2202.70),
    ("Ar Condicionado (instalação e manutenção)",    1960.00),
    ("Desinsetização (controle de pragas)",          1000.00),
]
build_detail_v7(
    "12", "Serviços", "contratados",
    5162.70, 5.3,
    "Serviços técnicos terceirizados: sistema de incêndio (Mar/Abr/Mai), ar condicionado (Fev/Jun) e desinsetização (Set).",
    SERV_TOTAL_MES,
    servicos_lanc
)

# 14 — Manutenção: subcategorias
manutencao_lanc = [
    ("Manutenção de Elevador (contrato preventivo mensal)",     2797.28),
    ("Manutenção de Segurança Eletrônica (sistema eletrônico)", 2312.64),
]
build_detail_v7(
    "13", "Contratos de", "manutenção",
    5109.92, 5.3,
    "Contratos preventivos mensais de elevador e sistema de segurança eletrônica. Dezembro sem cobrança por ajuste de competência.",
    MANUT_TOTAL_MES,
    manutencao_lanc
)

# 15 — Financeiras: mês a mês (tarifas) — só lançamentos com valor
fin_lanc = [(f"Tarifas Bancárias — {MESES_EXT[i]}/2025", FIN_TARIFAS[i]) for i in range(12)]
fin_lanc.append(("IRRF Poupança — Junho/2025", 2.30))
build_detail_v7(
    "14", "Despesas", "financeiras",
    2164.11, 2.2,
    "Tarifas bancárias mensais de manutenção de conta e uma retenção de IRRF sobre rendimentos de poupança em Jun.",
    FIN_TOTAL_MES,
    fin_lanc
)

# 16 — Retenções: meses com valor
ret_lanc = [(f"DAS Simples Nacional — {MESES_EXT[i]}/2025", RET_DAS[i]) for i in range(12) if RET_DAS[i] > 0]
build_detail_v7(
    "15", "Retenções e", "tributos",
    573.30, 0.6,
    "Retenção DAS Simples Nacional iniciada em Jun/2025 com parcela fixa mensal de R$ 81,90 até o fim do exercício.",
    RET_DAS,
    ret_lanc
)


# ============================================================
# SLIDE FINAL — ENCERRAMENTO com PAINEL ANTES x DEPOIS
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = C_BG_DARK
bg.line.fill.background(); bg.shadow.inherit = False

# Círculos decorativos
c_dec1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-2.5), Inches(6), Inches(6))
c_dec1.fill.solid(); c_dec1.fill.fore_color.rgb = RGBColor(0x16, 0x2E, 0x5E)
c_dec1.line.fill.background(); c_dec1.shadow.inherit = False

c_dec2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(5), Inches(5), Inches(5))
c_dec2.fill.solid(); c_dec2.fill.fore_color.rgb = RGBColor(0x13, 0x28, 0x54)
c_dec2.line.fill.background(); c_dec2.shadow.inherit = False

# Eyebrow
add_text_box(slide, Inches(0.5), Inches(0.4), Inches(10), Inches(0.3),
             "ENCERRAMENTO    •    BALANÇO DO EXERCÍCIO",
             size=11, bold=True, color=C_AMBER)
add_line(slide, Inches(0.5), Inches(0.78), Inches(1.0), Inches(0.78), C_AMBER, weight=1.8)

# Título
tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.33), Inches(0.9))
tf = tb.text_frame; tf.word_wrap = True; tf.margin_left = Emu(0); tf.margin_top = Emu(0)
p = tf.paragraphs[0]
r1 = p.add_run(); r1.text = "O ano em "
r1.font.name = FONT; r1.font.size = Pt(32); r1.font.bold = True; r1.font.color.rgb = C_WHITE
r2 = p.add_run(); r2.text = "números"
r2.font.name = FONT; r2.font.size = Pt(32); r2.font.bold = True; r2.font.color.rgb = C_AMBER

add_text_box(slide, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.35),
             "Como o patrimônio da Associação evoluiu ao longo de 2025",
             size=13, color=C_BLUE_PALE)

# ================================================
# MINI-GRÁFICO COMPARATIVO — CANTO SUPERIOR DIREITO
# ================================================
mini_x = Inches(8.5); mini_y = Inches(2.45)
mini_w = Inches(4.33); mini_h = Inches(2.2)
add_rounded_rect(slide, mini_x, mini_y, mini_w, mini_h, RGBColor(0x12, 0x22, 0x48), corner=0.05)

add_text_box(slide, mini_x + Inches(0.25), mini_y + Inches(0.15), mini_w - Inches(0.5), Inches(0.25),
             "COMPARATIVO DO SALDO", size=10, bold=True, color=C_AMBER)
add_text_box(slide, mini_x + Inches(0.25), mini_y + Inches(0.42), mini_w - Inches(0.5), Inches(0.22),
             "Início de 2025  vs  Fim de 2025", size=9, color=C_BLUE_PALE)

# Barras desenhadas manualmente para ter total controle visual
# Barras desenhadas manualmente
chart_base_y = mini_y + mini_h - Inches(0.35)  # linha de base das barras
max_bar_h = Inches(1.3)  # altura máxima = representa R$ 40.932,91
# Proporção: saldo inicial / saldo final
prop_inicial = SALDO_ANTERIOR / SALDO_FINAL  # ~0.514
h_inicial = Emu(int(max_bar_h.emu * prop_inicial))
h_final   = max_bar_h

bar_w = Inches(1.3)
gap_bars = Inches(0.4)
total_bars_w = bar_w * 2 + gap_bars
bar_start_x = mini_x + (mini_w - total_bars_w) / 2

# Barra ANTES (azul médio)
bar1_x = bar_start_x
bar1_y = chart_base_y - h_inicial
bar1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar1_x, bar1_y, bar_w, h_inicial)
bar1.fill.solid(); bar1.fill.fore_color.rgb = C_BLUE_MID
bar1.line.fill.background(); bar1.shadow.inherit = False

# Valor em cima da barra 1
add_text_box(slide, bar1_x - Inches(0.15), bar1_y - Inches(0.3),
             bar_w + Inches(0.3), Inches(0.25),
             "R$ 21.040", size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
# Label embaixo
add_text_box(slide, bar1_x - Inches(0.15), chart_base_y + Inches(0.05),
             bar_w + Inches(0.3), Inches(0.22),
             "INÍCIO 2025", size=8, bold=True, color=C_BLUE_PALE, align=PP_ALIGN.CENTER)

# Barra DEPOIS (laranja)
bar2_x = bar_start_x + bar_w + gap_bars
bar2_y = chart_base_y - h_final
bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar2_x, bar2_y, bar_w, h_final)
bar2.fill.solid(); bar2.fill.fore_color.rgb = C_AMBER
bar2.line.fill.background(); bar2.shadow.inherit = False

# Valor em cima da barra 2
add_text_box(slide, bar2_x - Inches(0.15), bar2_y - Inches(0.3),
             bar_w + Inches(0.3), Inches(0.25),
             "R$ 40.932", size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
# Label embaixo
add_text_box(slide, bar2_x - Inches(0.15), chart_base_y + Inches(0.05),
             bar_w + Inches(0.3), Inches(0.22),
             "FIM 2025", size=8, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)

# Linha base fina
add_line(slide,
         mini_x + Inches(0.3), chart_base_y,
         mini_x + mini_w - Inches(0.3), chart_base_y,
         C_BLUE_PALE, weight=0.5)

# ================================================
# 2 CARDS EMBAIXO do mini-comparativo
# ================================================
# Card 1 — Crescimento
m1_y = Inches(4.75); m1_h = Inches(1.05)
add_rounded_rect(slide, mini_x, m1_y, mini_w, m1_h, C_AMBER, corner=0.06)
add_text_box(slide, mini_x + Inches(0.25), m1_y + Inches(0.1), mini_w - Inches(0.5), Inches(0.22),
             "O CAIXA CRESCEU", size=9, bold=True, color=C_WHITE)
add_text_box(slide, mini_x + Inches(0.25), m1_y + Inches(0.3), mini_w - Inches(0.5), Inches(0.7),
             f"+{CRESCIMENTO_STR}%", size=38, bold=True, color=C_WHITE)

# Card 2 — Cobertura
m2_y = Inches(5.9); m2_h = Inches(1.05)
add_rounded_rect(slide, mini_x, m2_y, mini_w, m2_h, C_POSITIVE, corner=0.06)
add_text_box(slide, mini_x + Inches(0.25), m2_y + Inches(0.1), mini_w - Inches(0.5), Inches(0.22),
             f"RECEITA COBRIU {COBERTURA_STR}% DA DESPESA", size=9, bold=True, color=C_WHITE)
add_multi_run_text(
    slide, mini_x + Inches(0.25), m2_y + Inches(0.38),
    mini_w - Inches(0.5), Inches(0.55),
    [(f"{MESES_POSITIVOS} meses", 26, True, C_WHITE),
     (" com superávit", 13, False, C_WHITE)]
)

# ================================================
# GRÁFICO PRINCIPAL À ESQUERDA (com 2a linha do saldo inicial)
# ================================================
panel_x = Inches(0.5); panel_y = Inches(2.45)
panel_w = Inches(7.8); panel_h = Inches(4.35)
add_rounded_rect(slide, panel_x, panel_y, panel_w, panel_h, RGBColor(0x12, 0x22, 0x48), corner=0.03)

add_text_box(slide, panel_x + Inches(0.35), panel_y + Inches(0.25), panel_w - Inches(0.7), Inches(0.3),
             "EVOLUÇÃO DO PATRIMÔNIO EM CAIXA", size=11, bold=True, color=C_AMBER)
add_text_box(slide, panel_x + Inches(0.35), panel_y + Inches(0.55), panel_w - Inches(0.7), Inches(0.3),
             "De R$ 21.040,97 em janeiro para R$ 40.932,91 em dezembro",
             size=11, color=C_BLUE_PALE)

# Legenda manual
leg_y = panel_y + Inches(0.88)
# cor 1 - linha atual
add_rect(slide, panel_x + Inches(0.35), leg_y + Inches(0.08), Inches(0.2), Inches(0.06), C_AMBER)
add_text_box(slide, panel_x + Inches(0.6), leg_y, Inches(2.5), Inches(0.22),
             "Saldo ao longo de 2025", size=9, color=C_WHITE)
# cor 2 - linha de partida (base 2025)
add_rect(slide, panel_x + Inches(3.3), leg_y + Inches(0.08), Inches(0.2), Inches(0.06), C_BLUE_LIGHT)
add_text_box(slide, panel_x + Inches(3.55), leg_y, Inches(3.0), Inches(0.22),
             "Ponto de partida (início 2025)", size=9, color=C_WHITE)

# Gráfico com 2 séries
chart_data_final = CategoryChartData()
chart_data_final.categories = MESES
chart_data_final.add_series("Saldo em Caixa",        SALDO_FIM_MES)
chart_data_final.add_series("Ponto de Partida 2025", [SALDO_ANTERIOR]*12)

chart_shape_final = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, panel_x + Inches(0.2), panel_y + Inches(1.25),
    panel_w - Inches(0.4), panel_h - Inches(1.45), chart_data_final
)
chart_final = chart_shape_final.chart
chart_final.has_title = False
chart_final.has_legend = False

plot_f = chart_final.plots[0]
# Série principal — linha laranja grossa
s1 = plot_f.series[0]
s1.format.line.color.rgb = C_AMBER
s1.format.line.width = Pt(4)

# Série referência — linha azul clara fina tracejada
s2 = plot_f.series[1]
s2.format.line.color.rgb = C_BLUE_LIGHT
s2.format.line.width = Pt(1.5)
# Aplicar dash via XML
from pptx.oxml.ns import qn
spPr_s2 = s2.format.line._get_or_add_ln()
prstDash = spPr_s2.find(qn('a:prstDash'))
if prstDash is None:
    prstDash = spPr_s2.makeelement(qn('a:prstDash'), {'val': 'dash'})
    spPr_s2.append(prstDash)
else:
    prstDash.set('val', 'dash')

# Estilizar eixos
for ax in [chart_final.category_axis, chart_final.value_axis]:
    try:
        ax.tick_labels.font.size = Pt(9)
        ax.tick_labels.font.name = FONT
        ax.tick_labels.font.color.rgb = C_BLUE_PALE
    except Exception:
        pass

# ================================================
# TIRA COMPARATIVA INFERIOR (antes/depois)
# ================================================
comp_y = Inches(6.95)
comp_h = Inches(0.4)
tira_w = Inches(6.0)

add_rounded_rect(slide, Inches(0.5), comp_y, tira_w, comp_h, RGBColor(0x2A, 0x3A, 0x5E), corner=0.2)
add_multi_run_text(
    slide, Inches(0.75), comp_y + Inches(0.08), tira_w - Inches(0.5), Inches(0.28),
    [("INÍCIO DE 2025 ", 9, True, C_BLUE_PALE),
     ("•  Saldo em caixa: ", 10, False, C_WHITE),
     (fmt_brl(SALDO_ANTERIOR), 11, True, C_WHITE)]
)

add_rounded_rect(slide, Inches(6.83), comp_y, tira_w, comp_h, C_AMBER, corner=0.2)
add_multi_run_text(
    slide, Inches(7.08), comp_y + Inches(0.08), tira_w - Inches(0.5), Inches(0.28),
    [("FIM DE 2025 ", 9, True, C_WHITE),
     ("•  Saldo em caixa: ", 10, False, C_WHITE),
     (fmt_brl(SALDO_FINAL), 11, True, C_WHITE)]
)

output_path = "/home/claude/Prestacao_Contas_v10.pptx"
prs.save(output_path)
print(f"OK: {output_path}")
print(f"Total de slides: {len(prs.slides)}")
