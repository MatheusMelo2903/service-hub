"""
Gerador de Previsão Orçamentária Condominial — versão data-driven.

REGRA CRÍTICA: o fundo de reserva NÃO entra no rateio.
A taxa é calculada exclusivamente como: Despesa Operacional ÷ 12 ÷ Unidades Equivalentes.
O fundo aparece apenas como informação institucional na apresentação.

COMO USAR:
    python3 gerar_previsao.py <planilha.xlsx> [--condominio "Nome"] [--output saida.pptx]

A planilha deve ter as 4 abas padrão:
  - Reajustes
  - Previsao Anual
  - Previsao Mensal
  - Resumo Assembleia

Os reajustes são lidos DINAMICAMENTE da planilha, em qualquer combinação.
Pode haver:
  - Reajuste só em Funcionários
  - Reajuste em 6 categorias diferentes
  - Reajuste em itens específicos com % diferente da categoria
  - Reajuste em zero categorias (caso de manutenção da taxa)

A skill se adapta ao que estiver na planilha.

Slides 7+ (detalhamento) variam: uma para cada categoria com reajuste,
ordenados por peso decrescente. Se houver muitas categorias reajustadas,
gera múltiplos slides de detalhamento.
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


# ========================================================================
# PALETA
# ========================================================================
C_BG_DARK    = RGBColor(0x0A, 0x17, 0x33)
C_NAVY_DEEP  = RGBColor(0x0A, 0x24, 0x63)
C_NAVY       = RGBColor(0x14, 0x3A, 0x87)
C_BLUE       = RGBColor(0x1E, 0x5A, 0xA8)
C_BLUE_MID   = RGBColor(0x2E, 0x7B, 0xC7)
C_BLUE_LIGHT = RGBColor(0x52, 0x99, 0xDC)
C_BLUE_PALE  = RGBColor(0x7F, 0xB5, 0xE3)
C_GRAY_TEXT  = RGBColor(0x3E, 0x56, 0x76)
C_GRAY_MUTED = RGBColor(0x8B, 0x9A, 0xB8)
C_GRAY_BG    = RGBColor(0xF7, 0xF9, 0xFC)
C_GRAY_LINE  = RGBColor(0xE2, 0xE8, 0xF0)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_POSITIVE   = RGBColor(0x22, 0x8B, 0x54)
C_AMBER      = RGBColor(0xE8, 0x8B, 0x1A)
C_AMBER_DEEP = RGBColor(0xC4, 0x6E, 0x0A)
C_AMBER_LIGHT = RGBColor(0xFD, 0xE9, 0xCC)

CAT_COLORS = [
    RGBColor(0x1E, 0x73, 0x4A),
    RGBColor(0x35, 0x9E, 0x66),
    RGBColor(0x14, 0x3A, 0x87),
    RGBColor(0x1E, 0x5A, 0xA8),
    RGBColor(0x2E, 0x7B, 0xC7),
    RGBColor(0x52, 0x99, 0xDC),
    RGBColor(0x5B, 0x6A, 0x88),
    RGBColor(0x7F, 0x8F, 0xA8),
]

FONT = "Carlito"

# Descrições padrão por categoria — usadas nos cards de detalhamento e box.
# Se aparecer categoria nova fora desta lista, usa fallback genérico.
DESCRICOES_CATEGORIA = {
    "Despesa com Funcionários":
        "Folha de pagamento da equipe operacional do condomínio, com todos os "
        "encargos sociais, provisões trabalhistas e benefícios obrigatórios da categoria.",
    "Manutenção":
        "Manutenções recorrentes e preventivas dos equipamentos e áreas do "
        "condomínio: bombas, elevadores, jardinagem, segurança eletrônica e lazer.",
    "Consumo e Taxas":
        "Consumo de utilidades (energia, água, gás, telefonia) e taxas "
        "obrigatórias junto a órgãos públicos.",
    "Serviços":
        "Serviços técnicos contratados sob demanda: seguro, segurança e incêndio, "
        "elétrica, ar condicionado, reparos e melhorias.",
    "Despesa Administrativa":
        "Honorários da administradora, do síndico, custos jurídicos e cartoriais "
        "para a gestão formal do condomínio.",
    "Aquisição de Materiais":
        "Materiais consumíveis para operação contínua do condomínio: limpeza, "
        "manutenção, jardinagem, reposição de peças e atendimento ao sistema de incêndio.",
    "Despesas Financeiras":
        "Custos da operação bancária do condomínio: emissão de boletos, tarifas "
        "de manutenção de conta e taxas de retorno bancário.",
    "Equipamentos":
        "Aquisição de eletrodomésticos, ferramentas e equipamentos duráveis para "
        "uso na operação do condomínio.",
}

JUSTIFICATIVAS_CATEGORIA = {
    "Despesa com Funcionários":
        "Reajuste salarial e dissídio coletivo da categoria, com reflexo em encargos e provisões",
    "Manutenção":
        "Reajuste dos contratos de manutenção recorrente",
    "Consumo e Taxas":
        "Reajuste de tarifas de utilidades e taxas obrigatórias",
    "Serviços":
        "Reajuste dos contratos de serviços técnicos e seguro",
    "Despesa Administrativa":
        "Honorários da administradora, do síndico, cartório e advocatícios",
    "Aquisição de Materiais":
        "Inflação de materiais de construção, limpeza e manutenção",
    "Despesas Financeiras":
        "Reajuste das tarifas bancárias e custos de emissão de boletos",
    "Equipamentos":
        "Reposição e atualização de equipamentos da operação",
}


# ========================================================================
# LEITURA DA PLANILHA
# ========================================================================
ORDEM_CATEGORIAS_PADRAO = [
    "Despesas Financeiras",
    "Despesa com Funcionários",
    "Despesa Administrativa",
    "Consumo e Taxas",
    "Manutenção",
    "Aquisição de Materiais",
    "Equipamentos",
    "Serviços",
]


# ========================================================================
# CÁLCULOS
# ========================================================================
def calcular_taxas(dados):
    """Calcula taxas antes/depois sem incluir fundo no rateio."""
    base = dados["base_anual"]
    prev = dados["previsto"]
    ue = dados["unid_equiv"]
    fator = dados["fator_cobertura"]

    taxa_apto_antes = base / 12 / ue
    taxa_apto_depois = prev / 12 / ue
    taxa_cob_antes = taxa_apto_antes * fator
    taxa_cob_depois = taxa_apto_depois * fator
    diff_apto = taxa_apto_depois - taxa_apto_antes
    diff_cob = taxa_cob_depois - taxa_cob_antes
    impacto_pct = (prev - base) / base * 100 if base else 0.0

    return {
        "taxa_apto_antes": taxa_apto_antes,
        "taxa_apto_depois": taxa_apto_depois,
        "taxa_cob_antes": taxa_cob_antes,
        "taxa_cob_depois": taxa_cob_depois,
        "diff_apto": diff_apto,
        "diff_cob": diff_cob,
        "impacto_pct": impacto_pct,
    }


# ========================================================================
# HELPERS DE FORMATAÇÃO E SHAPES
# ========================================================================
def fmt_brl(v):
    s = f"{v:,.2f}"
    return "R$ " + s.replace(",", "X").replace(".", ",").replace("X", ".")


def fmt_pct(v):
    """Aceita 0.07 ou 7.0 e devolve '7%' ou '7,5%'."""
    if abs(v) < 1:
        v = v * 100
    s = f"{v:.2f}".replace(".", ",")
    if s.endswith(",00"):
        s = s[:-3]
    elif s.endswith("0"):
        s = s[:-1]
    return s + "%"


def fmt_pct_signed(v):
    s = fmt_pct(v)
    if v > 0 and not s.startswith("+"):
        s = "+" + s
    return s


def add_rect(slide, x, y, w, h, fill_color, radius=None):
    if radius is not None:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.adjustments[0] = radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=C_GRAY_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_line(slide, x1, y1, x2, y2, color, weight=1.0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


# ========================================================================
# CONSTRUÇÃO DOS SLIDES
# ========================================================================
class Builder:
    def __init__(self, dados, taxas, logo_path=None):
        self.dados = dados
        self.taxas = taxas
        self.cond_nome = dados["condominio"]
        self.logo_path = logo_path
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.BLANK = self.prs.slide_layouts[6]

        # Categorias ordenadas por peso decrescente
        prev_total = dados["previsto"]
        for c in dados["categorias"]:
            c["peso_pct"] = (c["previsto"] / prev_total * 100) if prev_total else 0.0
        self.cats_ordenadas = sorted(
            dados["categorias"], key=lambda c: c["previsto"], reverse=True
        )
        self.cats_com_reajuste = [c for c in self.cats_ordenadas if c["tem_reajuste"]]
        self.cats_sem_reajuste = [c for c in self.cats_ordenadas if not c["tem_reajuste"]]

    def footer(self, slide):
        add_text(
            slide, Inches(0.5), Inches(7.15), Inches(12.33), Inches(0.25),
            f"PREVISÃO ORÇAMENTÁRIA 2026    •    {self.cond_nome.upper()}    •    VIRTUAL SERVICE",
            size=9, color=C_GRAY_MUTED, align=PP_ALIGN.CENTER
        )

    def header(self, slide, number, section, title_main, title_accent=None, subtitle=None):
        add_text(slide, Inches(0.5), Inches(0.4), Inches(10), Inches(0.3),
                 f"{number}    {section}",
                 size=11, bold=True, color=C_BLUE_PALE)
        add_line(slide, Inches(0.5), Inches(0.78), Inches(1.0), Inches(0.78), C_BLUE_MID, weight=1.8)

        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.33), Inches(0.8))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(0); tf.margin_top = Emu(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run()
        r1.text = title_main
        r1.font.name = FONT; r1.font.size = Pt(32); r1.font.bold = True
        r1.font.color.rgb = C_NAVY
        if title_accent:
            r2 = p.add_run()
            r2.text = " " + title_accent
            r2.font.name = FONT; r2.font.size = Pt(32); r2.font.bold = True
            r2.font.color.rgb = C_BLUE_MID
        if subtitle:
            add_text(slide, Inches(0.5), Inches(1.85), Inches(12.33), Inches(0.4),
                     subtitle, size=14, color=C_GRAY_TEXT)

    # ============ SLIDE 1: CAPA ============
    def slide_capa(self):
        slide = self.prs.slides.add_slide(self.BLANK)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = C_BG_DARK
        bg.line.fill.background(); bg.shadow.inherit = False

        c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-2), Inches(6), Inches(6))
        c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x16, 0x2E, 0x5E)
        c1.line.fill.background(); c1.shadow.inherit = False

        c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(4.5), Inches(4.5), Inches(4.5))
        c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x13, 0x28, 0x54)
        c2.line.fill.background(); c2.shadow.inherit = False

        if self.logo_path and os.path.exists(self.logo_path):
            slide.shapes.add_picture(self.logo_path, Inches(10.8), Inches(0.55),
                                     width=Inches(2.0), height=Inches(0.66))

        add_text(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(0.4),
                 "PREVISÃO ORÇAMENTÁRIA", size=14, bold=True, color=C_BLUE_PALE)

        # Quebra "Condomínio Foo" em duas linhas
        nome = self.cond_nome
        if nome.lower().startswith("condomínio "):
            prefixo = "Condomínio"
            resto = nome[len("Condomínio "):]
        else:
            prefixo = nome
            resto = ""

        tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11), Inches(2.8))
        tf = tb.text_frame; tf.word_wrap = True; tf.margin_left = Emu(0)
        p1 = tf.paragraphs[0]
        r1 = p1.add_run(); r1.text = prefixo
        r1.font.name = FONT; r1.font.size = Pt(54); r1.font.bold = True
        r1.font.color.rgb = C_WHITE
        if resto:
            p2 = tf.add_paragraph()
            r2 = p2.add_run(); r2.text = resto
            r2.font.name = FONT; r2.font.size = Pt(54); r2.font.bold = True
            r2.font.color.rgb = C_BLUE_LIGHT

        add_line(slide, Inches(0.8), Inches(5.0), Inches(2.0), Inches(5.0), C_AMBER, weight=3)
        add_text(slide, Inches(0.8), Inches(5.25), Inches(10), Inches(0.4),
                 "Exercício 2026", size=20, bold=True, color=C_WHITE)
        add_text(slide, Inches(0.8), Inches(5.73), Inches(10), Inches(0.4),
                 "Apresentação à Assembleia de Moradores", size=13, color=C_BLUE_PALE)

        # Card impacto
        card_x = Inches(0.8); card_y = Inches(6.3); card_w = Inches(7.2); card_h = Inches(0.75)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h)
        card.adjustments[0] = 0.2
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x14, 0x2A, 0x55)
        card.line.color.rgb = C_AMBER; card.line.width = Pt(1.2)
        card.shadow.inherit = False

        impacto = self.taxas["impacto_pct"]
        diff_apto = self.taxas["diff_apto"]
        if abs(impacto) < 0.005:
            msg_esq = "Taxa condominial mantida"
        elif impacto < 1.5:
            msg_esq = f"Apenas {fmt_pct_signed(impacto)} na taxa"
        else:
            msg_esq = f"Reajuste de {fmt_pct_signed(impacto)} na taxa"

        add_text(slide, card_x + Inches(0.3), card_y + Inches(0.08), Inches(3.4), Inches(0.3),
                 "IMPACTO NA TAXA CONDOMINIAL", size=9, bold=True, color=C_AMBER)
        add_text(slide, card_x + Inches(0.3), card_y + Inches(0.35), Inches(3.4), Inches(0.4),
                 msg_esq, size=16, bold=True, color=C_WHITE)

        valor_dir = fmt_brl(diff_apto) if diff_apto >= 0 else "−" + fmt_brl(abs(diff_apto))
        add_text(slide, card_x + Inches(4.0), card_y + Inches(0.1), Inches(3.0), Inches(0.4),
                 valor_dir, size=24, bold=True, color=C_AMBER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, card_x + Inches(4.0), card_y + Inches(0.42), Inches(3.0), Inches(0.3),
                 "por mês por apartamento", size=10, color=C_BLUE_PALE)

    # ============ SLIDE 2: VISÃO GERAL ============
    def slide_visao_geral(self, slide_num_str="02"):
        slide = self.prs.slides.add_slide(self.BLANK)
        self.header(slide, slide_num_str, "VISÃO GERAL",
                    "A previsão orçamentária", "2026 em números",
                    subtitle="Síntese dos principais indicadores financeiros para o exercício de 2026")

        prev = self.dados["previsto"]
        fundo = self.dados["fundo_reserva"] or 0.0

        top_y = Inches(2.5); top_h = Inches(1.5); card_w = Inches(4.0); gap = Inches(0.17)
        cards_top = [
            ("DESPESA TOTAL PREVISTA", fmt_brl(prev), "Operação completa do condomínio em 2026", C_NAVY),
            ("FUNDO DE RESERVA", fmt_brl(fundo), "Reserva existente, intocável, fora do rateio", C_BLUE_MID),
            ("DESPESA MENSAL MÉDIA", fmt_brl(prev / 12), "Total previsto dividido por 12 meses", C_NAVY_DEEP),
        ]
        x = Inches(0.5)
        for label, val, desc, color in cards_top:
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_y, card_w, top_h)
            card.adjustments[0] = 0.08
            card.fill.solid(); card.fill.fore_color.rgb = color
            card.line.fill.background(); card.shadow.inherit = False
            add_text(slide, x + Inches(0.25), top_y + Inches(0.2), card_w - Inches(0.5), Inches(0.3),
                     label, size=10, bold=True, color=C_BLUE_PALE)
            add_text(slide, x + Inches(0.25), top_y + Inches(0.5), card_w - Inches(0.5), Inches(0.55),
                     val, size=22, bold=True, color=C_WHITE)
            add_text(slide, x + Inches(0.25), top_y + Inches(1.05), card_w - Inches(0.5), Inches(0.4),
                     desc, size=10, color=C_BLUE_PALE)
            x += card_w + gap

        # Taxas
        taxa_y = Inches(4.25); taxa_h = Inches(1.3); taxa_w = Inches(6.13)
        x = Inches(0.5)
        apt = self.dados["apartamentos"]
        cob = self.dados["coberturas"]
        fator = self.dados["fator_cobertura"]
        for label, val, desc in [
            ("TAXA MENSAL POR APARTAMENTO", fmt_brl(self.taxas["taxa_apto_depois"]),
             f"{apt} apartamentos no condomínio"),
            ("TAXA MENSAL POR COBERTURA", fmt_brl(self.taxas["taxa_cob_depois"]),
             f"{cob} coberturas, fator {str(fator).replace('.', ',')} sobre apartamento"),
        ]:
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, taxa_y, taxa_w, taxa_h)
            card.adjustments[0] = 0.06
            card.fill.solid(); card.fill.fore_color.rgb = C_GRAY_BG
            card.line.color.rgb = C_GRAY_LINE; card.line.width = Pt(0.75)
            card.shadow.inherit = False
            add_text(slide, x + Inches(0.3), taxa_y + Inches(0.2), taxa_w - Inches(0.6), Inches(0.3),
                     label, size=10, bold=True, color=C_BLUE_MID)
            add_text(slide, x + Inches(0.3), taxa_y + Inches(0.5), taxa_w - Inches(0.6), Inches(0.5),
                     val, size=28, bold=True, color=C_NAVY)
            add_text(slide, x + Inches(0.3), taxa_y + Inches(0.97), taxa_w - Inches(0.6), Inches(0.3),
                     desc, size=10, color=C_GRAY_TEXT)
            x += taxa_w + Inches(0.17)

        # Faixa âmbar
        faixa_y = Inches(5.85); faixa_h = Inches(1.0)
        faixa = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), faixa_y, Inches(12.33), faixa_h)
        faixa.adjustments[0] = 0.1
        faixa.fill.solid(); faixa.fill.fore_color.rgb = C_AMBER
        faixa.line.fill.background(); faixa.shadow.inherit = False

        total_cats = len(self.dados["categorias"])
        n_reaj = len(self.cats_com_reajuste)
        impacto_str = fmt_pct_signed(self.taxas["impacto_pct"])
        diff_str = fmt_brl(self.taxas["diff_apto"]) if self.taxas["diff_apto"] >= 0 else "−" + fmt_brl(abs(self.taxas["diff_apto"]))

        blocos = [
            (impacto_str, "IMPACTO TOTAL", "na taxa condominial em 2026"),
            (f"{n_reaj} de {total_cats}", "REAJUSTE FINANCEIRO", "categorias com reajuste aplicado"),
            (diff_str, "POR APARTAMENTO", "por mês a mais que em 2025"),
        ]
        bw = Inches(12.33) / 3
        bx = Inches(0.5)
        for v, lab, desc in blocos:
            add_text(slide, bx + Inches(0.2), faixa_y + Inches(0.13), bw - Inches(0.4), Inches(0.5),
                     v, size=24, bold=True, color=C_WHITE)
            add_text(slide, bx + Inches(0.2), faixa_y + Inches(0.55), bw - Inches(0.4), Inches(0.25),
                     lab, size=9, bold=True, color=RGBColor(0xFF, 0xE8, 0xC8))
            add_text(slide, bx + Inches(0.2), faixa_y + Inches(0.75), bw - Inches(0.4), Inches(0.25),
                     desc, size=10, color=C_WHITE)
            bx += bw

        self.footer(slide)

    # ============ SLIDE 3: COMPARATIVO ============
    def slide_comparativo(self, slide_num_str="03"):
        slide = self.prs.slides.add_slide(self.BLANK)
        self.header(slide, slide_num_str, "COMPARATIVO",
                    "Antes e depois", "dos reajustes 2026",
                    subtitle="Comparação entre a base realizada (sem reajuste) e a previsão orçamentária 2026")

        cy = Inches(2.7); ch = Inches(3.9); cw = Inches(6.0)

        # ANTES
        x1 = Inches(0.5)
        card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x1, cy, cw, ch)
        card1.adjustments[0] = 0.04
        card1.fill.solid(); card1.fill.fore_color.rgb = C_GRAY_BG
        card1.line.color.rgb = C_GRAY_LINE; card1.line.width = Pt(0.75)
        card1.shadow.inherit = False

        add_text(slide, x1 + Inches(0.35), cy + Inches(0.25), cw - Inches(0.7), Inches(0.3),
                 "SITUAÇÃO ATUAL", size=11, bold=True, color=C_GRAY_TEXT)
        add_text(slide, x1 + Inches(0.35), cy + Inches(0.55), cw - Inches(0.7), Inches(0.4),
                 "Base realizada em 2025", size=18, bold=True, color=C_NAVY)
        add_text(slide, x1 + Inches(0.35), cy + Inches(0.95), cw - Inches(0.7), Inches(0.3),
                 "Valores sem aplicação de reajustes", size=11, color=C_GRAY_TEXT)
        add_line(slide, x1 + Inches(0.35), cy + Inches(1.5), x1 + cw - Inches(0.35), cy + Inches(1.5),
                 C_GRAY_LINE, weight=0.75)
        add_text(slide, x1 + Inches(0.35), cy + Inches(1.65), cw - Inches(0.7), Inches(0.3),
                 "TAXA POR APARTAMENTO", size=9, bold=True, color=C_GRAY_TEXT)
        add_text(slide, x1 + Inches(0.35), cy + Inches(1.92), cw - Inches(0.7), Inches(0.5),
                 fmt_brl(self.taxas["taxa_apto_antes"]), size=28, bold=True, color=C_NAVY)
        add_text(slide, x1 + Inches(0.35), cy + Inches(2.55), cw - Inches(0.7), Inches(0.3),
                 "TAXA POR COBERTURA", size=9, bold=True, color=C_GRAY_TEXT)
        add_text(slide, x1 + Inches(0.35), cy + Inches(2.82), cw - Inches(0.7), Inches(0.5),
                 fmt_brl(self.taxas["taxa_cob_antes"]), size=28, bold=True, color=C_NAVY)
        add_text(slide, x1 + Inches(0.35), cy + Inches(3.45), cw - Inches(0.7), Inches(0.3),
                 f"Despesa base anual {fmt_brl(self.dados['base_anual'])}", size=10, color=C_GRAY_TEXT)

        # DEPOIS
        x2 = Inches(6.83)
        card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x2, cy, cw, ch)
        card2.adjustments[0] = 0.04
        card2.fill.solid(); card2.fill.fore_color.rgb = C_NAVY
        card2.line.fill.background(); card2.shadow.inherit = False

        add_text(slide, x2 + Inches(0.35), cy + Inches(0.25), cw - Inches(0.7), Inches(0.3),
                 "PREVISÃO 2026 COM REAJUSTES", size=11, bold=True, color=C_AMBER)
        add_text(slide, x2 + Inches(0.35), cy + Inches(0.55), cw - Inches(0.7), Inches(0.4),
                 "Orçamento previsto 2026", size=18, bold=True, color=C_WHITE)
        add_text(slide, x2 + Inches(0.35), cy + Inches(0.95), cw - Inches(0.7), Inches(0.3),
                 "Valores já com reajustes aplicados por categoria", size=11, color=C_BLUE_PALE)
        add_line(slide, x2 + Inches(0.35), cy + Inches(1.5), x2 + cw - Inches(0.35), cy + Inches(1.5),
                 C_BLUE, weight=0.75)
        add_text(slide, x2 + Inches(0.35), cy + Inches(1.65), cw - Inches(0.7), Inches(0.3),
                 "TAXA POR APARTAMENTO", size=9, bold=True, color=C_BLUE_PALE)
        add_text(slide, x2 + Inches(0.35), cy + Inches(1.92), cw - Inches(0.7), Inches(0.5),
                 fmt_brl(self.taxas["taxa_apto_depois"]), size=28, bold=True, color=C_WHITE)
        add_text(slide, x2 + Inches(0.35), cy + Inches(2.55), cw - Inches(0.7), Inches(0.3),
                 "TAXA POR COBERTURA", size=9, bold=True, color=C_BLUE_PALE)
        add_text(slide, x2 + Inches(0.35), cy + Inches(2.82), cw - Inches(0.7), Inches(0.5),
                 fmt_brl(self.taxas["taxa_cob_depois"]), size=28, bold=True, color=C_WHITE)
        add_text(slide, x2 + Inches(0.35), cy + Inches(3.45), cw - Inches(0.7), Inches(0.3),
                 f"Despesa total prevista {fmt_brl(self.dados['previsto'])}", size=10, color=C_BLUE_PALE)

        # Faixa âmbar inferior
        faixa_y = Inches(6.7); faixa_h = Inches(0.35)
        faixa = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), faixa_y, Inches(12.33), faixa_h)
        faixa.adjustments[0] = 0.15
        faixa.fill.solid(); faixa.fill.fore_color.rgb = C_AMBER
        faixa.line.fill.background(); faixa.shadow.inherit = False

        tb = slide.shapes.add_textbox(Inches(0.5), faixa_y, Inches(12.33), faixa_h)
        tf = tb.text_frame; tf.word_wrap = False
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        diff_apto_str = fmt_brl(self.taxas["diff_apto"]) if self.taxas["diff_apto"] >= 0 else "−" + fmt_brl(abs(self.taxas["diff_apto"]))
        diff_cob_str = fmt_brl(self.taxas["diff_cob"]) if self.taxas["diff_cob"] >= 0 else "−" + fmt_brl(abs(self.taxas["diff_cob"]))
        for txt, sz, bold in [
            ("IMPACTO POR APARTAMENTO  ", 11, True),
            (f"{'+' if self.taxas['diff_apto'] >= 0 else ''}{diff_apto_str}", 12, True),
            ("     •     IMPACTO POR COBERTURA  ", 11, True),
            (f"{'+' if self.taxas['diff_cob'] >= 0 else ''}{diff_cob_str}", 12, True),
        ]:
            r = p.add_run(); r.text = txt
            r.font.name = FONT; r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = C_WHITE

        add_text(slide, Inches(0.5), Inches(7.2), Inches(12.33), Inches(0.25),
                 f"PREVISÃO ORÇAMENTÁRIA 2026    •    {self.cond_nome.upper()}    •    VIRTUAL SERVICE",
                 size=9, color=C_GRAY_MUTED, align=PP_ALIGN.CENTER)

    # ============ SLIDE 4: REAJUSTES APLICADOS ============
    def slide_reajustes(self):
        slide = self.prs.slides.add_slide(self.BLANK)
        total_cats = len(self.dados["categorias"])
        n_reaj = len(self.cats_com_reajuste)
        if n_reaj == 0:
            subt = "Nenhuma categoria teve reajuste em 2026 — taxa mantida"
        else:
            subt = f"{n_reaj} das {total_cats} categorias receberam reajuste de inflação setorial em 2026"
        self.header(slide, "04", "REAJUSTES",
                    "Onde estão os reajustes", "para 2026", subtitle=subt)

        # 3 cards: as 3 categorias com reajuste de maior peso.
        # Se forem >3, mostramos as top 3 e citamos as demais no box inferior.
        # Se forem <3, mostramos só as que existem (sem cards vazios).
        cards_top = self.cats_com_reajuste[:3]

        cy = Inches(2.55); ch = Inches(3.35); gap = Inches(0.17)
        if len(cards_top) == 0:
            # Box explicativo único ocupando o lugar dos 3 cards
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), cy, Inches(12.33), ch)
            box.adjustments[0] = 0.05
            box.fill.solid(); box.fill.fore_color.rgb = C_GRAY_BG
            box.line.color.rgb = C_GRAY_LINE; box.line.width = Pt(0.75)
            box.shadow.inherit = False
            add_text(slide, Inches(0.8), cy + Inches(0.4), Inches(11.7), Inches(0.5),
                     "Sem reajustes aplicados para 2026", size=22, bold=True, color=C_NAVY)
            add_text(slide, Inches(0.8), cy + Inches(1.1), Inches(11.7), Inches(2.0),
                     "Todas as categorias foram mantidas no valor realizado em 2025. "
                     "A taxa condominial permanece igual ao exercício anterior.",
                     size=14, color=C_GRAY_TEXT)
        else:
            cw = Inches(12.33 / 3 - 0.12) if len(cards_top) == 3 else Inches(12.33 / len(cards_top) - 0.12)
            x = Inches(0.5)
            for cat in cards_top:
                self._card_reajuste(slide, x, cy, cw, ch, cat)
                x += cw + gap

        # Box inferior: categorias sem reajuste + (se houver) "outras com reajuste"
        box_y = Inches(6.05); box_h = Inches(0.95)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), box_y, Inches(12.33), box_h)
        box.adjustments[0] = 0.1
        box.fill.solid(); box.fill.fore_color.rgb = C_AMBER_LIGHT
        box.line.color.rgb = C_AMBER; box.line.width = Pt(0.75)
        box.shadow.inherit = False

        # Texto do box: muda conforme a situação
        sem_reaj_nomes = [c["nome"] for c in self.cats_sem_reajuste]
        extras_reaj = self.cats_com_reajuste[3:]  # as que não couberam nos 3 cards

        if extras_reaj:
            titulo_box = "OUTRAS CATEGORIAS COM REAJUSTE"
            partes = []
            for c in extras_reaj:
                partes.append(f"{c['nome']} ({fmt_pct_signed(c['reajuste_pct'])})")
            texto_box = "Também receberam reajuste em 2026: " + "; ".join(partes) + "."
            if sem_reaj_nomes:
                texto_box += " As demais categorias (" + ", ".join(sem_reaj_nomes) + ") foram mantidas no valor realizado em 2025."
        elif sem_reaj_nomes:
            titulo_box = "CATEGORIAS SEM REAJUSTE"
            peso_sem_reaj = sum(c["peso_pct"] for c in self.cats_sem_reajuste)
            texto_box = (
                f"{', '.join(sem_reaj_nomes)} foram mantidas no valor realizado em 2025. "
                f"Juntas representam cerca de {peso_sem_reaj:.0f}% do orçamento total, "
                "contribuindo para conter o impacto na taxa condominial."
            )
        else:
            titulo_box = "TODAS AS CATEGORIAS RECEBERAM REAJUSTE"
            texto_box = "Todas as categorias operacionais tiveram reajuste em 2026, conforme detalhado nos cards acima."

        add_text(slide, Inches(0.8), box_y + Inches(0.13), Inches(11.7), Inches(0.3),
                 titulo_box, size=10, bold=True, color=C_AMBER_DEEP)
        add_text(slide, Inches(0.8), box_y + Inches(0.4), Inches(11.7), Inches(0.55),
                 texto_box, size=10, color=C_GRAY_TEXT)

        self.footer(slide)

    def _card_reajuste(self, slide, x, cy, cw, ch, cat):
        """Renderiza um card vertical de reajuste de categoria."""
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, cw, ch)
        card.adjustments[0] = 0.05
        card.fill.solid(); card.fill.fore_color.rgb = C_GRAY_BG
        card.line.color.rgb = C_GRAY_LINE; card.line.width = Pt(0.75)
        card.shadow.inherit = False

        topo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, cw, Inches(0.95))
        topo.adjustments[0] = 0.1
        topo.fill.solid(); topo.fill.fore_color.rgb = C_NAVY
        topo.line.fill.background(); topo.shadow.inherit = False

        add_text(slide, x, cy + Inches(0.18), cw, Inches(0.6),
                 fmt_pct_signed(cat["reajuste_pct"]), size=36, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.3), cy + Inches(1.1), cw - Inches(0.6), Inches(0.5),
                 cat["nome"], size=15, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

        add_line(slide, x + Inches(0.8), cy + Inches(1.65), x + cw - Inches(0.8), cy + Inches(1.65),
                 C_GRAY_LINE, weight=0.5)

        add_text(slide, x + Inches(0.3), cy + Inches(1.8), cw - Inches(0.6), Inches(0.25),
                 "DE", size=9, bold=True, color=C_GRAY_MUTED, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.3), cy + Inches(2.0), cw - Inches(0.6), Inches(0.35),
                 fmt_brl(cat["base"]), size=14, bold=True, color=C_GRAY_TEXT, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.3), cy + Inches(2.35), cw - Inches(0.6), Inches(0.25),
                 "PARA", size=9, bold=True, color=C_AMBER_DEEP, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.3), cy + Inches(2.55), cw - Inches(0.6), Inches(0.4),
                 fmt_brl(cat["previsto"]), size=16, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

        just = JUSTIFICATIVAS_CATEGORIA.get(cat["nome"], "Reajuste de inflação setorial aplicado nesta categoria")
        add_text(slide, x + Inches(0.25), cy + Inches(2.97), cw - Inches(0.5), Inches(0.5),
                 just, size=9, color=C_GRAY_TEXT, align=PP_ALIGN.CENTER)

    # ============ SLIDE 5: ESTRUTURA DE DESPESAS ============
    def slide_estrutura(self):
        slide = self.prs.slides.add_slide(self.BLANK)
        total_cats = len(self.dados["categorias"])
        self.header(slide, "05", "ESTRUTURA",
                    "Como o orçamento", "está distribuído",
                    subtitle=f"Composição das {total_cats} categorias de despesa previstas para 2026")

        le_y = Inches(2.7)
        add_text(slide, Inches(0.5), le_y, Inches(5.5), Inches(0.3),
                 "DESPESA TOTAL PREVISTA", size=10, bold=True, color=C_BLUE_MID)
        add_text(slide, Inches(0.5), le_y + Inches(0.3), Inches(5.5), Inches(0.8),
                 fmt_brl(self.dados["previsto"]), size=38, bold=True, color=C_NAVY)
        add_text(slide, Inches(0.5), le_y + Inches(1.15), Inches(5.5), Inches(0.3),
                 "no orçamento 2026", size=12, color=C_GRAY_TEXT)
        add_text(slide, Inches(0.5), le_y + Inches(1.55), Inches(5.5), Inches(0.3),
                 f"Distribuído em {total_cats} categorias operacionais", size=11, color=C_GRAY_TEXT)

        # Insight dinâmico: 2 maiores categorias
        top2 = self.cats_ordenadas[:2]
        peso_top2 = sum(c["peso_pct"] for c in top2)
        nomes_top2 = " e ".join(c["nome"] for c in top2)
        insight_txt = (
            f"{nomes_top2} respondem juntas por mais de {peso_top2:.0f}% do orçamento, "
            "refletindo a operação contínua do condomínio."
        )

        insight_y = le_y + Inches(2.1)
        ins = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.5), insight_y, Inches(5.7), Inches(1.5))
        ins.adjustments[0] = 0.06
        ins.fill.solid(); ins.fill.fore_color.rgb = C_NAVY
        ins.line.fill.background(); ins.shadow.inherit = False
        add_text(slide, Inches(0.75), insight_y + Inches(0.2), Inches(5.2), Inches(0.3),
                 "INSIGHT", size=9, bold=True, color=C_AMBER)
        add_text(slide, Inches(0.75), insight_y + Inches(0.45), Inches(5.2), Inches(1.0),
                 insight_txt, size=11, color=C_WHITE)

        # Lista de categorias à direita
        ld_x = Inches(6.5); ld_y = Inches(2.5)
        line_h = 0.46
        for i, cat in enumerate(self.cats_ordenadas):
            y = ld_y + Inches(i * line_h)
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         ld_x, y + Inches(0.08), Inches(0.12), Inches(0.28))
            color_idx = i % len(CAT_COLORS)
            bar.fill.solid(); bar.fill.fore_color.rgb = CAT_COLORS[color_idx]
            bar.line.fill.background(); bar.shadow.inherit = False
            add_text(slide, ld_x + Inches(0.25), y, Inches(4.3), Inches(0.4),
                     cat["nome"], size=12, bold=True, color=C_NAVY, anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, ld_x + Inches(4.5), y, Inches(1.5), Inches(0.4),
                     fmt_brl(cat["previsto"]), size=11, bold=True, color=C_GRAY_TEXT,
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, ld_x + Inches(6.05), y, Inches(0.7), Inches(0.4),
                     f"{cat['peso_pct']:.1f}%".replace(".", ","), size=11, bold=True, color=C_BLUE_MID,
                     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

        self.footer(slide)

    # ============ SLIDE 6: DISTRIBUIÇÃO MENSAL ============
    def slide_mensal(self):
        slide = self.prs.slides.add_slide(self.BLANK)
        self.header(slide, "06", "MENSAL",
                    "Distribuição Mensal", "ao longo de 2026",
                    subtitle="Previsão mensal de despesas para 2026, com base no padrão realizado em 2025")

        previsao = self.dados["previsao_mensal"]
        meses = self.dados["meses"]
        media = sum(previsao) / 12
        max_v = max(previsao)
        min_v = min(previsao)
        max_idx = previsao.index(max_v)
        min_idx = previsao.index(min_v)

        chart_data = CategoryChartData()
        chart_data.categories = meses
        chart_data.add_series("Previsto", previsao)

        chart_x = Inches(0.5); chart_y = Inches(2.6); chart_w = Inches(8.5); chart_h = Inches(4.3)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, chart_x, chart_y, chart_w, chart_h, chart_data
        ).chart
        chart.has_title = False
        chart.has_legend = False

        plot = chart.plots[0]
        plot.gap_width = 90
        ser = plot.series[0]
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = C_BLUE_MID
        ser.format.line.fill.background()

        for ax in [chart.category_axis, chart.value_axis]:
            try:
                ax.tick_labels.font.size = Pt(10)
                ax.tick_labels.font.name = FONT
                ax.tick_labels.font.color.rgb = C_GRAY_TEXT
            except Exception:
                pass

        sx = Inches(9.3); sy = Inches(2.6); sw = Inches(3.5); sh = Inches(1.35); gap = Inches(0.1)
        side_cards = [
            ("MÊS DE MAIOR DESPESA", meses[max_idx], fmt_brl(max_v),
             "Pico de despesa previsto no ano", C_NAVY),
            ("MÉDIA MENSAL", "12 meses", fmt_brl(media),
             "Custo médio operacional mensal", C_BLUE_MID),
            ("MÊS DE MENOR DESPESA", meses[min_idx], fmt_brl(min_v),
             "Mês de menor custo previsto", C_POSITIVE),
        ]
        y = sy
        for label, mes, val, desc, color in side_cards:
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, y, sw, sh)
            card.adjustments[0] = 0.05
            card.fill.solid(); card.fill.fore_color.rgb = color
            card.line.fill.background(); card.shadow.inherit = False
            add_text(slide, sx + Inches(0.25), y + Inches(0.13), sw - Inches(0.5), Inches(0.25),
                     label, size=8.5, bold=True, color=C_AMBER_LIGHT)
            add_text(slide, sx + Inches(0.25), y + Inches(0.4), sw - Inches(0.5), Inches(0.3),
                     mes, size=11, bold=True, color=C_WHITE)
            add_text(slide, sx + Inches(0.25), y + Inches(0.65), sw - Inches(0.5), Inches(0.4),
                     val, size=17, bold=True, color=C_WHITE)
            add_text(slide, sx + Inches(0.25), y + Inches(1.05), sw - Inches(0.5), Inches(0.25),
                     desc, size=9, color=C_BLUE_PALE)
            y += sh + gap

        self.footer(slide)

    # ============ SLIDES 7+: DETALHAMENTO POR CATEGORIA REAJUSTADA ============
    def slide_detalhamento(self, slide_num_str, cat):
        slide = self.prs.slides.add_slide(self.BLANK)
        self.header(slide, slide_num_str, "DETALHAMENTO",
                    cat["nome"], "",
                    subtitle=f"Reajuste de {fmt_pct(cat['reajuste_pct'])} aplicado nesta categoria")

        # Card navy esquerdo
        card_x = Inches(0.5); card_y = Inches(2.55); card_w = Inches(3.8); card_h = Inches(4.5)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h)
        card.adjustments[0] = 0.04
        card.fill.solid(); card.fill.fore_color.rgb = C_NAVY
        card.line.fill.background(); card.shadow.inherit = False

        add_text(slide, card_x + Inches(0.3), card_y + Inches(0.25), card_w - Inches(0.6), Inches(0.3),
                 "CATEGORIA", size=10, bold=True, color=C_AMBER)
        add_text(slide, card_x + Inches(0.3), card_y + Inches(0.52), card_w - Inches(0.6), Inches(0.5),
                 cat["nome"], size=16, bold=True, color=C_WHITE)

        desc = DESCRICOES_CATEGORIA.get(
            cat["nome"],
            "Categoria de despesas operacionais do condomínio."
        )
        add_text(slide, card_x + Inches(0.3), card_y + Inches(1.2), card_w - Inches(0.6), Inches(0.9),
                 desc, size=10, color=C_BLUE_PALE)

        add_line(slide, card_x + Inches(0.3), card_y + Inches(1.95),
                 card_x + card_w - Inches(0.3), card_y + Inches(1.95), C_BLUE, weight=0.5)

        add_text(slide, card_x + Inches(0.3), card_y + Inches(2.1), card_w - Inches(0.6), Inches(0.25),
                 "BASE 2025", size=9, bold=True, color=C_BLUE_PALE)
        add_text(slide, card_x + Inches(0.3), card_y + Inches(2.35), card_w - Inches(0.6), Inches(0.4),
                 fmt_brl(cat["base"]), size=14, bold=True, color=C_WHITE)

        add_text(slide, card_x + Inches(0.3), card_y + Inches(2.85), card_w - Inches(0.6), Inches(0.25),
                 "PREVISTO 2026", size=9, bold=True, color=C_AMBER)
        add_text(slide, card_x + Inches(0.3), card_y + Inches(3.1), card_w - Inches(0.6), Inches(0.4),
                 fmt_brl(cat["previsto"]), size=14, bold=True, color=C_WHITE)

        comp_y = card_y + Inches(3.65)
        add_text(slide, card_x + Inches(0.3), comp_y, card_w - Inches(0.6), Inches(0.25),
                 "REAJUSTE APLICADO", size=9, bold=True, color=C_AMBER)
        diff = cat["previsto"] - cat["base"]
        add_text(slide, card_x + Inches(0.3), comp_y + Inches(0.22), card_w - Inches(0.6), Inches(0.4),
                 fmt_pct_signed(cat["reajuste_pct"]), size=18, bold=True, color=C_AMBER)
        if diff > 0:
            add_text(slide, card_x + Inches(0.3), comp_y + Inches(0.55), card_w - Inches(0.6), Inches(0.25),
                     f"(+{fmt_brl(diff)})", size=10, color=C_BLUE_PALE)

        # Tabela à direita: itens da categoria ordenados por valor previsto decrescente
        # 4 COLUNAS: BASE 2025 anual | MENSAL 2025 | MENSAL 2026 | PREVISTO 2026
        itens = sorted(cat["itens"], key=lambda i: i["previsto"], reverse=True)
        tx = Inches(4.55); ty = Inches(2.55); tw = Inches(8.28)
        head_h = Inches(0.4)
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, tx, ty, tw, head_h)
        head.fill.solid(); head.fill.fore_color.rgb = C_NAVY
        head.line.fill.background(); head.shadow.inherit = False

        col_item_w = Inches(3.0)
        col_w = Inches(1.32)
        add_text(slide, tx + Inches(0.15), ty, col_item_w, head_h,
                 "ITEM", size=9.5, bold=True, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, tx + Inches(0.15) + col_item_w, ty, col_w, head_h,
                 "BASE 2025", size=9.5, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(slide, tx + Inches(0.15) + col_item_w + col_w, ty, col_w, head_h,
                 "MENSAL 2025", size=9.5, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(slide, tx + Inches(0.15) + col_item_w + 2*col_w, ty, col_w, head_h,
                 "MENSAL 2026", size=9.5, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(slide, tx + Inches(0.15) + col_item_w + 3*col_w, ty,
                 col_w - Inches(0.15), head_h,
                 "PREVISTO 2026", size=9.5, bold=True, color=C_AMBER,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

        # Linhas: altura adaptativa conforme número de itens
        n_itens = len(itens)
        if n_itens <= 8:
            row_h = Inches(0.36); font_row = 9.5
        elif n_itens <= 12:
            row_h = Inches(0.28); font_row = 8.5
        elif n_itens <= 16:
            row_h = Inches(0.24); font_row = 8
        else:
            row_h = Inches(0.20); font_row = 7.5

        row_y = ty + head_h
        for i, item in enumerate(itens):
            if i % 2 == 0:
                bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, tx, row_y, tw, row_h)
                bg.fill.solid(); bg.fill.fore_color.rgb = C_GRAY_BG
                bg.line.fill.background(); bg.shadow.inherit = False
            add_text(slide, tx + Inches(0.15), row_y, col_item_w, row_h,
                     item["nome"], size=font_row, color=C_GRAY_TEXT, anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, tx + Inches(0.15) + col_item_w, row_y, col_w, row_h,
                     fmt_brl(item["base"]), size=font_row, color=C_GRAY_TEXT,
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
            add_text(slide, tx + Inches(0.15) + col_item_w + col_w, row_y, col_w, row_h,
                     fmt_brl(item["base"] / 12), size=font_row, color=C_GRAY_TEXT,
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
            add_text(slide, tx + Inches(0.15) + col_item_w + 2*col_w, row_y, col_w, row_h,
                     fmt_brl(item["previsto"] / 12), size=font_row, bold=True, color=C_NAVY,
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
            add_text(slide, tx + Inches(0.15) + col_item_w + 3*col_w, row_y,
                     col_w - Inches(0.15), row_h,
                     fmt_brl(item["previsto"]), size=font_row, bold=True, color=C_NAVY,
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
            row_y += row_h

        # Total
        total_y = row_y + Inches(0.1)
        # Verifica se cabe o total no slide (limite ~ 6.85")
        if total_y + Inches(0.45) > Inches(6.95):
            total_y = Inches(6.5)
        total = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, tx, total_y, tw, Inches(0.45))
        total.fill.solid(); total.fill.fore_color.rgb = C_AMBER
        total.line.fill.background(); total.shadow.inherit = False

        add_text(slide, tx + Inches(0.15), total_y, col_item_w, Inches(0.45),
                 "TOTAL DA CATEGORIA", size=10, bold=True, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, tx + Inches(0.15) + col_item_w, total_y, col_w, Inches(0.45),
                 fmt_brl(cat["base"]), size=10, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(slide, tx + Inches(0.15) + col_item_w + col_w, total_y, col_w, Inches(0.45),
                 fmt_brl(cat["base"] / 12), size=10, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(slide, tx + Inches(0.15) + col_item_w + 2*col_w, total_y, col_w, Inches(0.45),
                 fmt_brl(cat["previsto"] / 12), size=10, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(slide, tx + Inches(0.15) + col_item_w + 3*col_w, total_y,
                 col_w - Inches(0.15), Inches(0.45),
                 fmt_brl(cat["previsto"]), size=10, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

        self.footer(slide)

    # ============ SLIDE DETALHAMENTO SEM REAJUSTE (novo na v6) ============
    def slide_detalhamento_sem_reajuste(self, slide_num_str, cat):
        """Igual ao slide_detalhamento, mas com:
        - Tabela de 2 colunas (BASE 2025 anual e MENSAL 2025)
        - Sem cabeçalho âmbar PREVISTO 2026
        - Card lateral sem o bloco REAJUSTE
        - Tag âmbar 'MANTIDO EM 2026 — SEM REAJUSTE' no card lateral
        """
        slide = self.prs.slides.add_slide(self.BLANK)
        self.header(slide, slide_num_str, "DETALHAMENTO",
                    cat["nome"], "",
                    subtitle="Categoria mantida em 2026 com base no realizado de 2025")

        # Card navy esquerdo
        card_x = Inches(0.5); card_y = Inches(2.55); card_w = Inches(3.8); card_h = Inches(4.5)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h)
        card.adjustments[0] = 0.04
        card.fill.solid(); card.fill.fore_color.rgb = C_NAVY
        card.line.fill.background(); card.shadow.inherit = False

        add_text(slide, card_x + Inches(0.3), card_y + Inches(0.25), card_w - Inches(0.6), Inches(0.3),
                 "CATEGORIA", size=10, bold=True, color=C_AMBER)
        add_text(slide, card_x + Inches(0.3), card_y + Inches(0.52), card_w - Inches(0.6), Inches(0.5),
                 cat["nome"], size=16, bold=True, color=C_WHITE)

        desc = DESCRICOES_CATEGORIA.get(
            cat["nome"],
            "Categoria de despesas operacionais do condomínio."
        )
        add_text(slide, card_x + Inches(0.3), card_y + Inches(1.2), card_w - Inches(0.6), Inches(0.9),
                 desc, size=10, color=C_BLUE_PALE)

        add_line(slide, card_x + Inches(0.3), card_y + Inches(1.95),
                 card_x + card_w - Inches(0.3), card_y + Inches(1.95), C_BLUE, weight=0.5)

        # BASE 2025 e MENSAL 2025 (sem PREVISTO 2026)
        add_text(slide, card_x + Inches(0.3), card_y + Inches(2.1), card_w - Inches(0.6), Inches(0.25),
                 "BASE 2025", size=9, bold=True, color=C_BLUE_PALE)
        add_text(slide, card_x + Inches(0.3), card_y + Inches(2.35), card_w - Inches(0.6), Inches(0.4),
                 fmt_brl(cat["base"]), size=14, bold=True, color=C_WHITE)

        add_text(slide, card_x + Inches(0.3), card_y + Inches(2.85), card_w - Inches(0.6), Inches(0.25),
                 "MÉDIA MENSAL 2025", size=9, bold=True, color=C_BLUE_PALE)
        add_text(slide, card_x + Inches(0.3), card_y + Inches(3.1), card_w - Inches(0.6), Inches(0.4),
                 fmt_brl(cat["base"] / 12), size=14, bold=True, color=C_WHITE)

        # Tag inferior em âmbar
        tag_y = card_y + Inches(3.75)
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     card_x + Inches(0.3), tag_y,
                                     card_w - Inches(0.6), Inches(0.55))
        tag.adjustments[0] = 0.2
        tag.fill.solid(); tag.fill.fore_color.rgb = C_AMBER
        tag.line.fill.background(); tag.shadow.inherit = False
        add_text(slide, card_x + Inches(0.3), tag_y, card_w - Inches(0.6), Inches(0.55),
                 "MANTIDO EM 2026  •  SEM REAJUSTE",
                 size=10, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Tabela à direita: itens da categoria com 2 colunas (BASE ANUAL + MENSAL)
        itens = sorted(cat["itens"], key=lambda i: i["base"], reverse=True)
        tx = Inches(4.55); ty = Inches(2.55); tw = Inches(8.28)
        head_h = Inches(0.4)
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, tx, ty, tw, head_h)
        head.fill.solid(); head.fill.fore_color.rgb = C_NAVY
        head.line.fill.background(); head.shadow.inherit = False

        col_item_w = Inches(5.0)
        col_val_w = Inches(1.6)
        add_text(slide, tx + Inches(0.2), ty, col_item_w, head_h,
                 "ITEM", size=10, bold=True, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, tx + Inches(0.2) + col_item_w, ty, col_val_w, head_h,
                 "BASE 2025", size=10, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(slide, tx + Inches(0.2) + col_item_w + col_val_w, ty, col_val_w - Inches(0.2), head_h,
                 "MENSAL 2025", size=10, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

        n_itens = len(itens)
        if n_itens <= 8:
            row_h = Inches(0.36); font_row = 10
        elif n_itens <= 12:
            row_h = Inches(0.28); font_row = 9
        elif n_itens <= 16:
            row_h = Inches(0.24); font_row = 8.5
        else:
            row_h = Inches(0.20); font_row = 8

        row_y = ty + head_h
        for i, item in enumerate(itens):
            if i % 2 == 0:
                bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, tx, row_y, tw, row_h)
                bg.fill.solid(); bg.fill.fore_color.rgb = C_GRAY_BG
                bg.line.fill.background(); bg.shadow.inherit = False
            add_text(slide, tx + Inches(0.2), row_y, col_item_w, row_h,
                     item["nome"], size=font_row, color=C_GRAY_TEXT, anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, tx + Inches(0.2) + col_item_w, row_y, col_val_w, row_h,
                     fmt_brl(item["base"]), size=font_row, color=C_GRAY_TEXT,
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
            add_text(slide, tx + Inches(0.2) + col_item_w + col_val_w, row_y, col_val_w - Inches(0.2), row_h,
                     fmt_brl(item["base"] / 12), size=font_row, color=C_GRAY_TEXT,
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
            row_y += row_h

        # Total
        total_y = row_y + Inches(0.1)
        if total_y + Inches(0.45) > Inches(6.95):
            total_y = Inches(6.5)
        total = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, tx, total_y, tw, Inches(0.45))
        total.fill.solid(); total.fill.fore_color.rgb = C_AMBER
        total.line.fill.background(); total.shadow.inherit = False
        add_text(slide, tx + Inches(0.2), total_y, col_item_w, Inches(0.45),
                 "TOTAL DA CATEGORIA", size=10, bold=True, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, tx + Inches(0.2) + col_item_w, total_y, col_val_w, Inches(0.45),
                 fmt_brl(cat["base"]), size=11, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(slide, tx + Inches(0.2) + col_item_w + col_val_w, total_y, col_val_w - Inches(0.2), Inches(0.45),
                 fmt_brl(cat["base"] / 12), size=11, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

        self.footer(slide)

    # ============ SLIDE PANORAMA UNIFICADO (novo na v6) ============
    def slide_panorama_unificado(self, slide_num_str="03"):
        """Slide unificado 2025 vs 2026:
        - Header
        - Dois cards de total lado a lado (cinza 2025 + navy 2026)
        - Tabela com 5 colunas: CATEGORIA, MENSAL 25, ANUAL 25, MENSAL 26, ANUAL 26
        - Linha TOTAL âmbar
        - Card de insight com as 2 maiores categorias
        """
        slide = self.prs.slides.add_slide(self.BLANK)
        self.header(slide, slide_num_str, "PANORAMA",
                    "Quanto o condomínio paga hoje", "vs vai pagar em 2026",
                    subtitle="Comparação direta entre o realizado de 2025 e a proposta para 2026")

        base = self.dados["base_anual"]
        prev = self.dados["previsto"]

        # ===== Dois cards de total no topo (esquerda) =====
        card_y = Inches(2.55); card_h = Inches(1.5); card_w = Inches(3.35); gap = Inches(0.15)

        # Card 2025 (cinza claro)
        x1 = Inches(0.5)
        c1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x1, card_y, card_w, card_h)
        c1.adjustments[0] = 0.08
        c1.fill.solid(); c1.fill.fore_color.rgb = C_GRAY_BG
        c1.line.color.rgb = C_GRAY_LINE; c1.line.width = Pt(0.75)
        c1.shadow.inherit = False
        add_text(slide, x1 + Inches(0.25), card_y + Inches(0.2), card_w - Inches(0.5), Inches(0.3),
                 "REALIZADO 2025", size=10, bold=True, color=C_GRAY_TEXT)
        add_text(slide, x1 + Inches(0.25), card_y + Inches(0.5), card_w - Inches(0.5), Inches(0.55),
                 fmt_brl(base), size=22, bold=True, color=C_NAVY)
        add_text(slide, x1 + Inches(0.25), card_y + Inches(1.05), card_w - Inches(0.5), Inches(0.4),
                 f"média mensal {fmt_brl(base / 12)}", size=10, color=C_GRAY_TEXT)

        # Card 2026 (navy)
        x2 = x1 + card_w + gap
        c2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x2, card_y, card_w, card_h)
        c2.adjustments[0] = 0.08
        c2.fill.solid(); c2.fill.fore_color.rgb = C_NAVY
        c2.line.fill.background(); c2.shadow.inherit = False
        add_text(slide, x2 + Inches(0.25), card_y + Inches(0.2), card_w - Inches(0.5), Inches(0.3),
                 "PREVISTO 2026", size=10, bold=True, color=C_AMBER)
        add_text(slide, x2 + Inches(0.25), card_y + Inches(0.5), card_w - Inches(0.5), Inches(0.55),
                 fmt_brl(prev), size=22, bold=True, color=C_WHITE)
        add_text(slide, x2 + Inches(0.25), card_y + Inches(1.05), card_w - Inches(0.5), Inches(0.4),
                 f"média mensal {fmt_brl(prev / 12)}", size=10, color=C_BLUE_PALE)

        # ===== Tabela à direita: categorias =====
        # Deslocada um pouco mais pra esquerda e com colunas anuais maiores
        tx = Inches(7.4); ty = Inches(2.55); tw = Inches(5.43)
        head_h = Inches(0.4)
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, tx, ty, tw, head_h)
        head.fill.solid(); head.fill.fore_color.rgb = C_NAVY
        head.line.fill.background(); head.shadow.inherit = False

        # 5 colunas: categoria + 4 valores. Anuais maiores que mensais.
        col_cat_w = Inches(1.85)
        col_mens_w = Inches(0.78)
        col_anual_w = Inches(1.01)
        add_text(slide, tx + Inches(0.1), ty, col_cat_w, head_h,
                 "CATEGORIA", size=8.5, bold=True, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, tx + col_cat_w, ty, col_mens_w, head_h,
                 "MENS 25", size=8, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(slide, tx + col_cat_w + col_mens_w, ty, col_anual_w, head_h,
                 "ANUAL 25", size=8, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(slide, tx + col_cat_w + col_mens_w + col_anual_w, ty, col_mens_w, head_h,
                 "MENS 26", size=8, bold=True, color=C_AMBER,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(slide, tx + col_cat_w + 2*col_mens_w + col_anual_w, ty,
                 col_anual_w - Inches(0.1), head_h,
                 "ANUAL 26", size=8, bold=True, color=C_AMBER,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

        # Linhas: 8 categorias + 1 total = 9 linhas
        n_cats = len(self.cats_ordenadas)
        if n_cats <= 8:
            row_h = Inches(0.32); font_row = 8
        else:
            row_h = Inches(0.28); font_row = 7.5

        row_y = ty + head_h
        for i, cat in enumerate(self.cats_ordenadas):
            if i % 2 == 0:
                bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, tx, row_y, tw, row_h)
                bg.fill.solid(); bg.fill.fore_color.rgb = C_GRAY_BG
                bg.line.fill.background(); bg.shadow.inherit = False
            # Nome
            add_text(slide, tx + Inches(0.1), row_y, col_cat_w, row_h,
                     cat["nome"], size=font_row, color=C_GRAY_TEXT, anchor=MSO_ANCHOR.MIDDLE)
            # Mensal 25
            add_text(slide, tx + col_cat_w, row_y, col_mens_w, row_h,
                     fmt_brl(cat["base"] / 12), size=font_row, color=C_GRAY_TEXT,
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
            # Anual 25
            add_text(slide, tx + col_cat_w + col_mens_w, row_y, col_anual_w, row_h,
                     fmt_brl(cat["base"]), size=font_row, color=C_GRAY_TEXT,
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
            # Mensal 26 (navy bold)
            add_text(slide, tx + col_cat_w + col_mens_w + col_anual_w, row_y, col_mens_w, row_h,
                     fmt_brl(cat["previsto"] / 12), size=font_row, bold=True, color=C_NAVY,
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
            # Anual 26 (navy bold)
            add_text(slide, tx + col_cat_w + 2*col_mens_w + col_anual_w, row_y,
                     col_anual_w - Inches(0.1), row_h,
                     fmt_brl(cat["previsto"]), size=font_row, bold=True, color=C_NAVY,
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
            row_y += row_h

        # Linha total âmbar
        total = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, tx, row_y, tw, Inches(0.4))
        total.fill.solid(); total.fill.fore_color.rgb = C_AMBER
        total.line.fill.background(); total.shadow.inherit = False
        add_text(slide, tx + Inches(0.1), row_y, col_cat_w, Inches(0.4),
                 "TOTAL", size=9, bold=True, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, tx + col_cat_w, row_y, col_mens_w, Inches(0.4),
                 fmt_brl(base / 12), size=8, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(slide, tx + col_cat_w + col_mens_w, row_y, col_anual_w, Inches(0.4),
                 fmt_brl(base), size=8, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(slide, tx + col_cat_w + col_mens_w + col_anual_w, row_y, col_mens_w, Inches(0.4),
                 fmt_brl(prev / 12), size=8, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(slide, tx + col_cat_w + 2*col_mens_w + col_anual_w, row_y,
                 col_anual_w - Inches(0.1), Inches(0.4),
                 fmt_brl(prev), size=8, bold=True, color=C_WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

        # ===== Card de insight à esquerda, abaixo dos dois cards de total =====
        if len(self.cats_ordenadas) >= 2:
            top2 = self.cats_ordenadas[:2]
            soma_pct = sum(c["peso_pct"] for c in top2)
            nomes = " e ".join(c["nome"] for c in top2)
            insight_y = Inches(4.35); insight_w = Inches(6.75); insight_h = Inches(1.95)
            ins = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.5), insight_y, insight_w, insight_h)
            ins.adjustments[0] = 0.06
            ins.fill.solid(); ins.fill.fore_color.rgb = C_NAVY
            ins.line.fill.background(); ins.shadow.inherit = False

            add_text(slide, Inches(0.8), insight_y + Inches(0.25), insight_w - Inches(0.6), Inches(0.3),
                     "INSIGHT", size=10, bold=True, color=C_AMBER)
            tb = slide.shapes.add_textbox(Inches(0.8), insight_y + Inches(0.6),
                                          insight_w - Inches(0.6), insight_h - Inches(0.8))
            tf = tb.text_frame; tf.word_wrap = True
            tf.margin_left = Emu(0); tf.margin_top = Emu(0)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            for txt, sz, bold, color in [
                (nomes, 14, True, C_WHITE),
                (" respondem juntas por ", 14, False, C_BLUE_PALE),
                (f"{soma_pct:.0f}%".replace(".", ","), 14, True, C_AMBER),
                (" do orçamento, refletindo a operação contínua do condomínio.",
                 14, False, C_BLUE_PALE),
            ]:
                r = p.add_run(); r.text = txt
                r.font.name = FONT; r.font.size = Pt(sz); r.font.bold = bold
                r.font.color.rgb = color

        self.footer(slide)

    # ============ SLIDE ENCERRAMENTO "MUITO OBRIGADO" (novo na v6) ============
    def slide_encerramento_obrigado(self):
        """Slide final institucional:
        - Fundo navy escuro (como a capa)
        - Logo Grupo Service grande e centralizado
        - 'Muito obrigado' centralizado
        - Tagline 'Qualidade. Excelência. Transparência.'
        - Círculos decorativos como na capa
        - Sem footer numerado
        """
        slide = self.prs.slides.add_slide(self.BLANK)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = C_BG_DARK
        bg.line.fill.background(); bg.shadow.inherit = False

        # Círculos decorativos
        c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(-3), Inches(-2), Inches(6), Inches(6))
        c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x16, 0x2E, 0x5E)
        c1.line.fill.background(); c1.shadow.inherit = False

        c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(10.5), Inches(4.5), Inches(4.5), Inches(4.5))
        c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x13, 0x28, 0x54)
        c2.line.fill.background(); c2.shadow.inherit = False

        # Logo centralizado no topo (grande)
        if self.logo_path and os.path.exists(self.logo_path):
            # Logo grande: 3.6" largura. Centralizar horizontalmente.
            logo_w = Inches(3.6); logo_h = Inches(1.18)
            logo_x = (self.prs.slide_width - logo_w) / 2
            slide.shapes.add_picture(self.logo_path, logo_x, Inches(1.3),
                                     width=logo_w, height=logo_h)

        # Linha âmbar decorativa
        line_w = Inches(1.5)
        line_x = (self.prs.slide_width - line_w) / 2
        add_line(slide, line_x, Inches(3.05), line_x + line_w, Inches(3.05),
                 C_AMBER, weight=3)

        # "Muito obrigado" centralizado
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.35),
                                      self.prs.slide_width - Inches(1.0), Inches(1.2))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(0); tf.margin_top = Emu(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = "Muito obrigado"
        r.font.name = FONT; r.font.size = Pt(60); r.font.bold = True
        r.font.color.rgb = C_WHITE

        # Tagline
        add_text(slide, Inches(0.5), Inches(4.85),
                 self.prs.slide_width - Inches(1.0), Inches(0.5),
                 "Qualidade.  Excelência.  Transparência.",
                 size=18, color=C_BLUE_LIGHT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Subtítulo discreto: nome do condomínio
        add_text(slide, Inches(0.5), Inches(5.7),
                 self.prs.slide_width - Inches(1.0), Inches(0.4),
                 f"Previsão Orçamentária 2026  •  {self.cond_nome}",
                 size=11, color=C_BLUE_PALE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        add_text(slide, Inches(0.5), Inches(6.15),
                 self.prs.slide_width - Inches(1.0), Inches(0.4),
                 "Apresentação à Assembleia de Moradores",
                 size=10, color=C_GRAY_MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ============ SLIDE METODOLOGIA ============
    def slide_metodologia(self, slide_num_str):
        slide = self.prs.slides.add_slide(self.BLANK)
        self.header(slide, slide_num_str, "METODOLOGIA",
                    "Como o cálculo", "foi feito",
                    subtitle="Transparência no processo de definição da taxa condominial para 2026")

        cy = Inches(2.6); ch = Inches(2.6); cw = Inches(3.0); gap = Inches(0.11)
        total_cats = len(self.dados["categorias"])
        n_reaj = len(self.cats_com_reajuste)
        fundo_str = fmt_brl(self.dados["fundo_reserva"] or 0)
        ue_str = str(self.dados["unid_equiv"]).replace(".", ",")
        if ue_str.endswith(",0"):
            ue_str = ue_str[:-2]

        metodologia_data = [
            ("1", "BASE", "Realizado de 12 meses do condomínio, extraído do sistema Superlógica",
             fmt_brl(self.dados["base_anual"]), "Despesa total realizada"),
            ("2", "REAJUSTES",
             "Reajustes aplicados apenas em categorias com alta volatilidade ou inflação acumulada",
             f"{n_reaj} de {total_cats}", "Categorias reajustadas"),
            ("3", "FUNDO DE RESERVA",
             "Reserva já existente do condomínio, intocável e destinada exclusivamente a emergências e obras imprevistas (informacional)",
             fundo_str, "Não compõe o rateio"),
            ("4", "RATEIO",
             f"Divisão entre apartamentos ({self.dados['apartamentos']}) e coberturas ({self.dados['coberturas']}), com fator de {str(self.dados['fator_cobertura']).replace('.', ',')} sobre apartamento",
             ue_str, "Unidades equivalentes"),
        ]
        x = Inches(0.5)
        for num, titulo, desc, valor, valor_label in metodologia_data:
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, cw, ch)
            card.adjustments[0] = 0.05
            card.fill.solid(); card.fill.fore_color.rgb = C_GRAY_BG
            card.line.color.rgb = C_GRAY_LINE; card.line.width = Pt(0.75)
            card.shadow.inherit = False

            circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                          x + Inches(0.25), cy + Inches(0.2),
                                          Inches(0.55), Inches(0.55))
            circ.fill.solid(); circ.fill.fore_color.rgb = C_NAVY
            circ.line.fill.background(); circ.shadow.inherit = False
            add_text(slide, x + Inches(0.25), cy + Inches(0.22), Inches(0.55), Inches(0.55),
                     num, size=18, bold=True, color=C_AMBER,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

            add_text(slide, x + Inches(0.9), cy + Inches(0.3), cw - Inches(1.1), Inches(0.4),
                     titulo, size=12, bold=True, color=C_NAVY, anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, x + Inches(0.3), cy + Inches(0.95), cw - Inches(0.6), Inches(1.0),
                     desc, size=10, color=C_GRAY_TEXT)
            add_line(slide, x + Inches(0.3), cy + Inches(1.95),
                     x + cw - Inches(0.3), cy + Inches(1.95), C_GRAY_LINE, weight=0.5)
            add_text(slide, x + Inches(0.3), cy + Inches(2.08), cw - Inches(0.6), Inches(0.3),
                     valor, size=15, bold=True, color=C_BLUE_MID)
            add_text(slide, x + Inches(0.3), cy + Inches(2.35), cw - Inches(0.6), Inches(0.2),
                     valor_label, size=9, color=C_GRAY_MUTED)
            x += cw + gap

        # Fórmula
        form_y = Inches(5.4); form_h = Inches(1.5)
        form = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.5), form_y, Inches(12.33), form_h)
        form.adjustments[0] = 0.06
        form.fill.solid(); form.fill.fore_color.rgb = C_NAVY_DEEP
        form.line.fill.background(); form.shadow.inherit = False

        add_text(slide, Inches(0.8), form_y + Inches(0.2), Inches(11.7), Inches(0.3),
                 "FÓRMULA DA TAXA POR APARTAMENTO", size=10, bold=True, color=C_AMBER)

        tb = slide.shapes.add_textbox(Inches(0.8), form_y + Inches(0.55), Inches(11.7), Inches(0.5))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(0); tf.margin_top = Emu(0)
        p = tf.paragraphs[0]
        for txt, sz, bold, color in [
            ("Despesa Operacional Prevista", 14, True, C_WHITE),
            (f" ÷ 12 meses ÷ {ue_str} unidades = ", 14, False, C_BLUE_PALE),
            (fmt_brl(self.taxas["taxa_apto_depois"]), 16, True, C_AMBER),
        ]:
            r = p.add_run(); r.text = txt
            r.font.name = FONT; r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = color

        add_text(slide, Inches(0.8), form_y + Inches(1.05), Inches(11.7), Inches(0.35),
                 f"{fmt_brl(self.dados['previsto'])} ÷ 12 ÷ {ue_str} = {fmt_brl(self.taxas['taxa_apto_depois'])} por mês          •          Fundo de reserva é informacional, não compõe o rateio",
                 size=11, color=C_BLUE_PALE)

        self.footer(slide)

    # ============ SLIDE ENCERRAMENTO ============
    def slide_encerramento(self):
        slide = self.prs.slides.add_slide(self.BLANK)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = C_BG_DARK
        bg.line.fill.background(); bg.shadow.inherit = False

        c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(5), Inches(5))
        c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x13, 0x28, 0x54)
        c1.line.fill.background(); c1.shadow.inherit = False

        add_text(slide, Inches(0.5), Inches(0.5), Inches(10), Inches(0.3),
                 "ENCERRAMENTO", size=11, bold=True, color=C_BLUE_PALE)
        add_line(slide, Inches(0.5), Inches(0.88), Inches(1.0), Inches(0.88), C_AMBER, weight=2)

        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(1.5))
        tf = tb.text_frame; tf.word_wrap = True; tf.margin_left = Emu(0); tf.margin_top = Emu(0)
        p1 = tf.paragraphs[0]
        r1 = p1.add_run(); r1.text = "A previsão"
        r1.font.name = FONT; r1.font.size = Pt(32); r1.font.bold = True; r1.font.color.rgb = C_WHITE
        r1b = p1.add_run(); r1b.text = " em síntese"
        r1b.font.name = FONT; r1b.font.size = Pt(32); r1b.font.bold = True; r1b.font.color.rgb = C_BLUE_LIGHT

        add_text(slide, Inches(0.5), Inches(2.05), Inches(12), Inches(0.4),
                 f"Síntese da previsão orçamentária 2026 do {self.cond_nome}",
                 size=13, color=C_BLUE_PALE)

        comp_y = Inches(2.9)
        add_text(slide, Inches(0.5), comp_y, Inches(7.5), Inches(0.3),
                 "TAXA MENSAL POR APARTAMENTO", size=10, bold=True, color=C_AMBER)

        bar_w_total = Inches(6.5)
        antes = self.taxas["taxa_apto_antes"]
        depois = self.taxas["taxa_apto_depois"]
        ratio = (antes / depois) if depois > 0 else 1.0

        b1_y = comp_y + Inches(0.5)
        b1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), b1_y, bar_w_total * ratio, Inches(0.5))
        b1.fill.solid(); b1.fill.fore_color.rgb = C_GRAY_MUTED
        b1.line.fill.background(); b1.shadow.inherit = False
        add_text(slide, Inches(0.5), b1_y + Inches(0.55), Inches(2), Inches(0.25),
                 "ANTES", size=9, bold=True, color=C_BLUE_PALE)
        add_text(slide, Inches(0.5) + bar_w_total * ratio + Inches(0.15), b1_y + Inches(0.08),
                 Inches(1.5), Inches(0.35), fmt_brl(antes), size=12, bold=True, color=C_WHITE)

        b2_y = b1_y + Inches(1.0)
        b2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), b2_y, bar_w_total, Inches(0.5))
        b2.fill.solid(); b2.fill.fore_color.rgb = C_AMBER
        b2.line.fill.background(); b2.shadow.inherit = False
        add_text(slide, Inches(0.5), b2_y + Inches(0.55), Inches(2), Inches(0.25),
                 "DEPOIS", size=9, bold=True, color=C_AMBER)
        add_text(slide, Inches(0.5) + bar_w_total + Inches(0.15), b2_y + Inches(0.08),
                 Inches(1.5), Inches(0.35), fmt_brl(depois), size=12, bold=True, color=C_WHITE)

        # Cards laterais
        n_reaj = len(self.cats_com_reajuste)
        fundo_pct_str = f"{int(self.dados['fundo_pct'] * 100)}% Fundo"
        diff_apto = self.taxas["diff_apto"]
        impacto = self.taxas["impacto_pct"]
        diff_apto_str = fmt_brl(diff_apto) if diff_apto >= 0 else "−" + fmt_brl(abs(diff_apto))

        rx = Inches(8.7); ry = Inches(2.9); rw = Inches(4.3); rh = Inches(1.1); gap = Inches(0.15)
        right_cards = [
            (fmt_pct_signed(impacto),
             f"Aumento total de {diff_apto_str} por mês no apartamento" if diff_apto >= 0 else "Taxa mantida no exercício 2026",
             C_AMBER),
            (f"{n_reaj} categoria{'s' if n_reaj != 1 else ''}",
             "Receberam reajustes de inflação setorial" if n_reaj > 0 else "Sem reajustes em 2026",
             C_BLUE_LIGHT),
            (fundo_pct_str, "Reserva para imprevistos e obras", C_AMBER_LIGHT),
        ]
        y = ry
        for val, desc, color in right_cards:
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, y, rw, rh)
            card.adjustments[0] = 0.08
            card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x14, 0x2A, 0x55)
            card.line.color.rgb = color; card.line.width = Pt(1.2)
            card.shadow.inherit = False
            add_text(slide, rx + Inches(0.25), y + Inches(0.15), rw - Inches(0.5), Inches(0.45),
                     val, size=18, bold=True, color=color)
            add_text(slide, rx + Inches(0.25), y + Inches(0.6), rw - Inches(0.5), Inches(0.45),
                     desc, size=10, color=C_WHITE)
            y += rh + gap

        # Tira antes/depois
        tira_y = Inches(6.5); tira_h = Inches(0.7)
        tira = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.5), tira_y, Inches(12.33), tira_h)
        tira.adjustments[0] = 0.1
        tira.fill.solid(); tira.fill.fore_color.rgb = RGBColor(0x14, 0x2A, 0x55)
        tira.line.color.rgb = C_AMBER; tira.line.width = Pt(1)
        tira.shadow.inherit = False

        tb = slide.shapes.add_textbox(Inches(0.5), tira_y, Inches(12.33), tira_h)
        tf = tb.text_frame; tf.word_wrap = False
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        for txt, sz, bold, color in [
            ("APARTAMENTO  ", 11, True, C_BLUE_PALE),
            (fmt_brl(self.taxas["taxa_apto_antes"]), 14, True, C_GRAY_MUTED),
            ("  →  ", 14, True, C_AMBER),
            (fmt_brl(self.taxas["taxa_apto_depois"]), 16, True, C_WHITE),
            ("        COBERTURA  ", 11, True, C_BLUE_PALE),
            (fmt_brl(self.taxas["taxa_cob_antes"]), 14, True, C_GRAY_MUTED),
            ("  →  ", 14, True, C_AMBER),
            (fmt_brl(self.taxas["taxa_cob_depois"]), 16, True, C_WHITE),
        ]:
            r = p.add_run(); r.text = txt
            r.font.name = FONT; r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color

        add_text(slide, Inches(0.5), Inches(7.2), Inches(12.33), Inches(0.25),
                 f"PREVISÃO ORÇAMENTÁRIA 2026    •    {self.cond_nome.upper()}    •    VIRTUAL SERVICE",
                 size=9, color=C_GRAY_MUTED, align=PP_ALIGN.CENTER)

    # ============ MONTA TUDO (v6) ============
    def build(self):
        """Ordem v6:
        01  Capa
        02  Metodologia
        03  Panorama unificado 2025 vs 2026
        04+ Detalhamentos COM reajuste (por peso decrescente)
        ..  Detalhamentos SEM reajuste (por peso decrescente)
        N-2 Comparativo
        N-1 Visão geral
        N   Encerramento "Muito obrigado"
        """
        numero = 1

        # 01 — Capa (não usa número no header)
        self.slide_capa()
        numero += 1

        # 02 — Metodologia
        self.slide_metodologia(f"{numero:02d}")
        numero += 1

        # 03 — Panorama unificado
        self.slide_panorama_unificado(f"{numero:02d}")
        numero += 1

        # Detalhamentos COM reajuste (por peso decrescente)
        for cat in self.cats_com_reajuste:
            self.slide_detalhamento(f"{numero:02d}", cat)
            numero += 1

        # Detalhamentos SEM reajuste (por peso decrescente)
        for cat in self.cats_sem_reajuste:
            self.slide_detalhamento_sem_reajuste(f"{numero:02d}", cat)
            numero += 1

        # Penúltimo - 1: Comparativo
        self.slide_comparativo(f"{numero:02d}")
        numero += 1

        # Penúltimo: Visão geral
        self.slide_visao_geral(f"{numero:02d}")
        numero += 1

        # Último: Encerramento "Muito obrigado" (sem número)
        self.slide_encerramento_obrigado()

        return self.prs


# ========================================================================
# MAIN
# ========================================================================

__all__ = ['Builder', 'calcular_taxas', 'FONT']
